import copy
from pptx import Presentation

def update_ppt_with_data(ppt_file: str, data: Dict, output_file: str, rally_report_generator: RallyReportGenerator, start_date, end_date: str = None, teams=[], num_milestone_duplicates: int = 1):
    try:
        prs = Presentation(ppt_file)
        milestone_slide_index = -1

        # Find the index of the milestone slide
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    try:
                        first_row_cells = [cell.text.strip() for cell in shape.table.rows[0].cells]
                        if first_row_cells[0] == "Completed Milestones" and first_row_cells[1] == "Incomplete Milestones":
                            milestone_slide_index = i
                            print(f"Milestone table found on slide index: {i}")
                            break  # Found it, break inner loop
                    except Exception:
                        continue  # Table might be empty or malformed
            if milestone_slide_index != -1:
                break  # Found it, break outer loop

        if milestone_slide_index != -1:
            original_slide = prs.slides[milestone_slide_index]

            # Insert the specified number of duplicate slides
            for i in range(num_milestone_duplicates):
                new_slide_index = milestone_slide_index + 1 + i
                blank_slide_layout = prs.slide_layouts[6]  # Or original_slide.slide_layout
                new_slide = prs.slides.insert(new_slide_index, blank_slide_layout)

                # Copy all shapes from the original slide to the new slide
                for original_shape in original_slide.shapes:
                    el = original_shape.element
                    newel = copy.deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

            print(f"Milestone slide at index {milestone_slide_index} duplicated {num_milestone_duplicates} times.")

        # Reset the milestone update counter
        rally_report_generator.milestone_update_counter = 0

        # Update the content of the tables (including the duplicated ones)
        slide_index_counter = 0
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    is_milestone_table = False
                    try:
                        first_row_cells = [cell.text.strip() for cell in shape.table.rows[0].cells]
                        if first_row_cells[0] == "Completed Milestones" and first_row_cells[1] == "Incomplete Milestones":
                            is_milestone_table = True
                    except Exception:
                        pass

                    if is_milestone_table:
                        print(f"Updating milestone table on slide index: {i}, counter: {rally_report_generator.milestone_update_counter}")
                        _update_table(shape.table, data, rally_report_generator)
                    else:
                        # Update other tables if needed
                        pass

        prs.save(output_file)
        return True

    except Exception as e:
        print(f"Error updating PPT with data: {e}")
        return False
