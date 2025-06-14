from io import BytesIO
from docx.shared import RGBColor as doc_color
from docx import Document
import copy
import streamlit as st
from datetime import datetime, timedelta
import logging
import os
import sys
import requests
from collections import defaultdict
from pptx.dml.color import RGBColor 
from pptx import Presentation
from pptx.util import Pt
from typing import Dict, List, Optional, Tuple, Any
from pptx.enum.text import MSO_ANCHOR

# Configuration 
API_KEY = os.environ.get("RALLY_API_KEY", "_GlhfJLBwSBWjJhQOR01I18U2Synak0upT01INX0s") #token is invalid
PPT_TEMPLATE_PATH = os.environ.get("PPT_TEMPLATE_PATH", "template-Mark.pptx")
LOG_LEVEL = os.environ.get("LOG_LEVEL", "INFO").upper()
TEAM_OWNERS_EMIS = ["lakshminarayana nainaru", "Govindarajan M", "Ravi Ranjan"]

# Logging 
logging.basicConfig(level=LOG_LEVEL, format="%(asctime)s - %(levelname)s - %(message)s")

# RallyReportGenerator class and related functions
class RallyReportGenerator:
    """Fetches data from Rally and formats it for a report."""
    def __init__(self, api_key: str = API_KEY):
        self.teams: Dict[str, Dict[str, Dict[str, Any]]] = {}
        self.status_index: Dict[str, Dict[str, Dict[str, Any]]] = {}
        self.headers: Dict[str, str] = {"zsessionid": api_key}
        self.milestone_update_counter = 0
        self.milestone_counter = 1
        logging.debug("RallyReportGenerator initialized.")

    def _make_request(self, url: str, params: Optional[Dict] = None) -> Optional[Dict]:
        try:
            response = requests.get(url, headers=self.headers, params=params)
            response.raise_for_status()
            return response.json()
        except requests.exceptions.RequestException as e:
            logging.error(f"Request error: {e} for URL: {url}")
            return None

    def get_rally_projects(self):
        url = "https://rally1.rallydev.com/slm/webservice/v2.0/project"
        response = self._make_request(url)
        if not response or not response.get('QueryResult', {}).get('Results'):
            logging.error("Failed to fetch projects from Rally.")
            return None
        projects = response.get('QueryResult', {}).get('Results',[])
        return [project.get('_refObjectName', '') for project in projects]

    def add_story(self, team_name: str, story_id: str, story_details: Dict) -> None:
        self.teams.setdefault(team_name, {}).setdefault(story_details['status'], {})[story_id] = story_details
        self.status_index.setdefault(story_details['status'], {}).setdefault(team_name, {})[story_id] = story_details
        logging.debug(f"Story added: {team_name} - {story_id} - {story_details['status']}")

    def get_flex_resource_info(self, team: Optional[List[str]] = None, start_date: Optional[str] = None) -> Dict[str, str]:
        team = ["megha.chakraborty@aig.com", "Aarthi.Panneerselvam@aig.com", "vukyam.srisravya@aig.com"]
        start_date = start_date or datetime.today().strftime('%Y-%m-%d')
        url = f"https://rally1.rallydev.com/slm/webservice/v2.0/iteration?query=((StartDate <= \"{start_date}\") and (EndDate >= \"{start_date}\"))"
        flex_data: Dict[str, str] = defaultdict(str)
        counters: Dict[str, int] = defaultdict(int)
        response_data = self._make_request(url)
        if response_data:
            for iteration in response_data.get('QueryResult', {}).get('Results', []):
                for owner in team:
                    query_params = {"query": f'(Owner = "{owner}")'}
                    workproducts = self._make_request(iteration.get('_ref', "") + "/workproducts", params=query_params)
                    if workproducts:
                        for item in workproducts.get('QueryResult', {}).get('Results', []):
                            name = item.get('Name')
                            bau_team = item.get('c_DataBAUTeam', '') or item.get('Project', {}).get('_refObjectName')
                            owner_name = f"{item.get('Owner', {}).get('_refObjectName')}"
                            task_estimate = item.get('TaskEstimateTotal', '')
                            counters[owner_name] += 1
                            flex_data[owner_name] += f"{counters[owner_name]}. {name} ({task_estimate})\n"
                    else:
                        logging.warning(f"No work product for {owner}")
        return flex_data

    def fetch_iteration_dates(self, start_date: Optional[str] = None, next: bool = False) -> Any:
        start_date = start_date or datetime.today().strftime('%Y-%m-%d')
        url = f"https://rally1.rallydev.com/slm/webservice/v2.0/iteration?query=((StartDate <= \"{start_date}\") and (EndDate >= \"{start_date}\"))"
        response = self._make_request(url)
        if not response or not response.get('QueryResult', {}).get('Results'):
            return None
        response_data = self._make_request(response.get('QueryResult', {}).get('Results', [])[0].get('_ref', ''))
        if not response_data:
            return None
        data = response_data.get('Iteration', {})
        if next:
            return (datetime.strptime(data.get('EndDate', ''), "%Y-%m-%dT%H:%M:%S.%fZ") + timedelta(days=1)).strftime("%Y-%m-%d")
        return {
            "Start": datetime.strptime(data.get('StartDate', ''), "%Y-%m-%dT%H:%M:%S.%fZ").strftime('%Y-%m-%d'),
            "end": (datetime.strptime(data.get('EndDate', ""), "%Y-%m-%dT%H:%M:%S.%fZ") - timedelta(days=1)).strftime('%Y-%m-%d'),
            "iteration": data.get('_refObjectName', ''),
            "state": data.get('State', '')
        }

    def get_stories_by_status(self, status: Optional[str] = None, team: Optional[str] = None) -> Dict:
        if team and status:
            team_stories = self.teams.get(team, {})
            return [item['title'].strip() for item in team_stories.get(status, {}).values()]
        if team:
            return self.teams.get(team, {})
        if not status:
            return self.teams
        return self.status_index.get(status, {})

    def _process_workproduct(self, project_name: str, workproduct: Dict, nextIter: Optional[bool] = False, teams: List = []) -> Tuple[Optional[str], str]:
        
        task_status = workproduct.get('ScheduleState', "")
        name = workproduct.get('Name', "")
        display_color = workproduct.get('DisplayColor', "")
        team = workproduct.get('c_DataBAUTeam', "") 
        imp_data_key = None
        story_key = f'{project_name}{len(self.teams.get(project_name, {}).get(task_status, {}))}'

        if project_name in teams:
            if project_name == "RDM":
                imp_data_key = 'RDM-midsprint' if display_color == '#fce205' else 'RDM'
                if not nextIter:
                    self.add_story('RDM', story_key, {'title': name, 'status': task_status})
            elif project_name == "Data BAU":
                if team == "KPI":
                    imp_data_key = 'KPI-midsprint' if display_color == '#fce205' else 'KPI'
                    if not nextIter:
                        self.add_story('KPI', story_key, {'title': name, 'status': task_status})
                elif team == "EDW":
                    imp_data_key = 'EDW-midsprint' if display_color == '#fce205' else 'EDW'
                    if not nextIter:
                        self.add_story('EDW', story_key, {'title': name, 'status': task_status})
                elif team == "Trade Credit":
                    imp_data_key = 'Trade Credit-midsprint' if display_color == '#fce205' else 'Trade Credit'
                    if not nextIter:
                        self.add_story('Trade Credit', story_key, {'title': name, 'status': task_status})
            elif project_name == "CDL":
                imp_data_key = 'CDL-midsprint' if display_color == '#fce205' else 'CDL'
                if not nextIter:
                    self.add_story('CDL', story_key, {'title': name, 'status': task_status})
            elif project_name == "SCUP NA Datamart - KPI":
                imp_data_key = 'SCUP NA-midsprint' if display_color == '#fce205' else 'SCUP NA'
                if not nextIter:
                    self.add_story('SCUP NA', story_key, {'title': name, 'status': task_status})
            elif project_name == "CDH":
                imp_data_key = 'CDH-midsprint' if display_color == '#fce205' else 'CDH'
                if not nextIter:
                    self.add_story('CDH', story_key, {'title': name, 'status': task_status})
            elif project_name == "ADB":
                owner = workproduct.get('Owner', {}).get('_refObjectName', "")
                if display_color == '#fce205':
                    if team == "EMIS":
                        if owner in TEAM_OWNERS_EMIS:
                            imp_data_key = 'EMIS Backend-midsprint'
                            if not nextIter:
                                self.add_story('EMIS Backend', f'EMIS Backend{len(self.teams.get("EMIS Backend", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                        else:
                            imp_data_key = 'EMIS UI-midsprint'
                            if not nextIter:
                                self.add_story('EMIS UI', f'EMIS UI{len(self.teams.get("EMIS UI", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                    elif team == "CBIP":
                        imp_data_key = 'CBIP-midsprint'
                        if not nextIter:
                            self.add_story('CBIP', f'CBIP{len(self.teams.get("CBIP", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                    elif team == "D&B":
                        imp_data_key = 'D&B-midsprint'
                        if not nextIter:
                            self.add_story('D&B', f'D&B{len(self.teams.get("D&B", {}).get(task_status, {}))}', {'title': name, 'status': task_status})

                else:
                    if team == "EMIS":
                        if owner in TEAM_OWNERS_EMIS:
                            imp_data_key = 'EMIS Backend'
                            if not nextIter:
                                self.add_story('EMIS Backend', f'EMIS Backend{len(self.teams.get("EMIS Backend", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                        else:
                            imp_data_key = 'EMIS UI'
                            if not nextIter:
                                self.add_story('EMIS UI', f'EMIS UI{len(self.teams.get("EMIS UI", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                    elif team == "CBIP":
                        imp_data_key = 'CBIP'
                        if not nextIter:
                            self.add_story('CBIP', f'CBIP{len(self.teams.get("CBIP", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                    elif team == "D&B":
                        imp_data_key = 'D&B'
                        if not nextIter:
                            self.add_story('D&B', f'CBIP{len(self.teams.get("D&B", {}).get(task_status, {}))}', {'title': name, 'status': task_status})

        return imp_data_key, name

    def fetch_iteration_data(self, start_date: Optional[str] = None, next: Optional[bool] = False, teams: List = []) -> Optional[Dict[str, str]]:
        start_date = start_date or datetime.today().strftime('%Y-%m-%d')
        url = f"https://rally1.rallydev.com/slm/webservice/v2.0/iteration?query=((StartDate <= \"{start_date}\") and (EndDate >= \"{start_date}\"))"
        response_data = self._make_request(url)
        if not response_data:
            return None
        imp_data: Dict[str, str] = defaultdict(str, {
            'CBIP': '', 'EMIS Backend': '', 'EMIS UI': '',
            'CBIP-midsprint': '', 'EMIS Backend-midsprint': '',
            'EMIS UI-midsprint': '',
            'RDM': '', 'RDM-midsprint': '',
            'D&B': '', 'D&B-midsprint': '',
            'EDW': '', 'EDW-midsprint': '',
            'KPI': '', 'KPI-midsprint': '',
            'CDH': '', 'CDH-midsprint': '',
            'CDL': '', 'CDL-midsprint': '',
            'Trade Credit': '', 'Trade Credit-midsprint': '',
            'SCUP NA': '', 'SCUP NA-midsprint': ''
        })
        counters: Dict[str, int] = defaultdict(int)

        for iteration in response_data.get('QueryResult', {}).get('Results', []):
            iter_details = self._make_request(iteration.get('_ref', ""))
            if not iter_details:
                continue

            iter_data = iter_details.get('Iteration', {})
            logging.debug(f"Processing iteration: {iter_data.get('Name', [])}, {iter_data.get('Project', {}).get('_refObjectName', [])}")
            project_name = iter_data.get('Project', {}).get('_refObjectName', [])

            if project_name in teams:
                workproducts_ref = iter_data.get('WorkProducts', {}).get('_ref', "")
                if workproducts_ref:
                    workproducts_data = self._make_request(workproducts_ref)
                    if not workproducts_data:
                        continue
                    workproducts = workproducts_data.get('QueryResult', {}).get('Results', [])
                    for workproduct in workproducts:
                        imp_data_key, name = self._process_workproduct(project_name, workproduct, next, teams)
                        if imp_data_key:
                            counters[imp_data_key] += 1
                            imp_data[imp_data_key] += f"{counters[imp_data_key]}. {name}\n"
        return imp_data

    def get_all_data(self, start_date: Optional[str] = None, teams = [], end_date: Optional[str] = None) -> Dict[str, Any]:
        start_date = start_date or datetime.today().strftime('%Y-%m-%d')
        if end_date:
            return {
            "iter_dates": self.fetch_iteration_dates(start_date=start_date),
            "milestone_data" : self.milestonedetails(start_date=start_date, teams=teams, end_date=end_date, split=False),
        }

        return {
            "iteration_data": self.fetch_iteration_data(start_date=start_date, teams=teams),
            "iter_dates": self.fetch_iteration_dates(start_date=start_date),
            #"flex_resource_data": self.get_flex_resource_info(start_date=start_date),
            "deployed_data": self.get_stories_by_status(status="Deployed"),
            "milestone_data" : self.milestonedetails(start_date=start_date, teams=teams),
            "next_iteration_data": self.fetch_iteration_data(start_date=self.fetch_iteration_dates(start_date=start_date, next=True), teams=teams, next=True)
        }

    def milestonedetails(self, start_date: Optional[str] = None, end_date: Optional[str] = None, teams = [], split = True ):
        if start_date and end_date:
            if not validate_date_format(start_date):
                logging.error("Invalid date format used for milestone start_date param. Please use YYYY-MM-DD format.")
                sys.exit(1)

            if not validate_date_format(end_date):
                logging.error("Invalid date format used for milestone end_date param. Please use YYYY-MM-DD format.")
                sys.exit(1)

            start = start_date
            end = end_date
        else:
            iter_dates = self.fetch_iteration_dates(start_date=start_date)
            start = iter_dates['Start']
            end = iter_dates['end']

        projects = [{'RDM': 370746842872}, {'ADB': 81259836048}, {'CDH': 501753244176}, {'SCUP NA Datamart - KPI': 812858540333}, {'Data BAU': 343403251580}]
        
        filtered_projects = [next(iter(project.values())) for project in [project for project in projects if next(iter(project)) in teams]]
      
        data = defaultdict(list, { 
            'Active': [],
            'Inactive': []
        })
        for i in filtered_projects:
            url = f'https://rally1.rallydev.com/slm/webservice/v2.0/milestone?fetch=Artifacts,Name,DisplayColor,FormattedID,TargetDate&query=(((Projects contains "/project/{i}") AND (TargetDate >= {start})) AND (TargetDate <= {end}))&start=1&pagesize=25&order=TargetDate DESC'
            response_data = self._make_request(url)
            if not response_data:
                return None
            for milestone in response_data.get('QueryResult', {}).get('Results', []):
                Milestone_name = milestone.get('_refObjectName','')
                Milestone_id = milestone.get('FormattedID','')
                Milestone_color = milestone.get('DisplayColor','')
                artifact_details = self._make_request(milestone.get('Artifacts', {}).get('_ref',''))
                artifact_data = [f"{result['_refObjectName']} - {result['FormattedID']}" for result in artifact_details.get('QueryResult', {}).get('Results',[])]
                if Milestone_color == "#107c1e":
                    data['Active'] += [{"Milestone":Milestone_name + ' - [' + Milestone_id + ']', "len":len(artifact_data), "us":artifact_data}]
                else:
                    data['Inactive'] += [{"Milestone":Milestone_name + ' - [' + Milestone_id + ']', "len":len(artifact_data), "us":artifact_data}]
        
        if not split:
            print(data)
            return data
        active_milestones = self.distribute_paragraphs_to_slides(milestone_info=data['Active'], lines_per_slide=20)
        inactive_milestones = self.distribute_paragraphs_to_slides(milestone_info=data['Inactive'], lines_per_slide=20)

        data = {
            "Active" :  active_milestones,
            "Inactive" : inactive_milestones
        }   
        print(data) 
        return data

    def distribute_paragraphs_to_slides(self, milestone_info, lines_per_slide):
        paragraph_lengths = [i.get('len') for i in milestone_info]
        paragraph_indices = sorted(range(len(paragraph_lengths)), key=lambda i: paragraph_lengths[i], reverse=True)
        slides = []
        remaining_lines = []
        for index in paragraph_indices:
            length = paragraph_lengths[index]
            placed = False
            for i in range(len(slides)):
                if remaining_lines[i] >= length:
                    slides[i].append(index)
                    remaining_lines[i] -= length
                    placed = True
                    break
            if not placed:
                slides.append([index])
                remaining_lines.append(lines_per_slide - length)
        # Actual sorting of Milestone objects       
        result = []
        
        for indices in slides:
            extracted_objects = [milestone_info[index] for index in indices if index < len(milestone_info)]
            result.append(extracted_objects)
        return result
    
    def process_tasks(self, task, team):
        def task_order(x):
            if len(x) < 4:
                return False
            y = x.split(". ", 1)[1].strip()
            return not (y not in self.get_stories_by_status(team=team, status="In-Progress") and
                        y not in self.get_stories_by_status(team=team, status="Idea") and
                        y not in self.get_stories_by_status(team=team, status="Defined"))
        sorted_tasks = sorted(task, key=task_order)
        tasks = [task for task in sorted_tasks if task]
        tasks_string = "\n".join(task.split(". ", 1)[1] for task in tasks)
        return tasks_string

def get_formatted_titles(data: Dict, team: str) -> str:
    titles = [value["title"] for value in data.get(team, {}).values()]
    return "\n".join(f"{i+1}. {title}" for i, title in enumerate(titles))

def _update_table(table: Any, new_planned_user_stories: Dict, rally_report_generator: RallyReportGenerator) -> None:
    if len(table.rows[0].cells) > 2:
        if table.rows[0].cells[0].text.strip() == "Application" and table.rows[0].cells[1].text.strip() == "Planned user stories" and table.rows[0].cells[2].text.strip() == "Mid Sprint user stories":
            for row in table.rows:
                app_name = row.cells[0].text.strip()
                if app_name in new_planned_user_stories['iteration_data']:
                    text_frame = row.cells[1].text_frame
                    text_frame.clear()
                    text = rally_report_generator.process_tasks(new_planned_user_stories['iteration_data'][app_name].split('\n'), team=app_name)
                    text_frame.text = text if text else "None"
                    for i, paragraph in enumerate(text_frame.paragraphs):
                        paragraph.font.size = Pt(12)
                        if len(paragraph.text) > 4:
                            if paragraph.text.strip() not in rally_report_generator.get_stories_by_status(team=app_name, status="In-Progress"):
                                if paragraph.text.strip() not in rally_report_generator.get_stories_by_status(team=app_name, status="Idea"):
                                    if paragraph.text.strip() not in rally_report_generator.get_stories_by_status(team=app_name, status="Defined"):
                                        if paragraph.text.strip().startswith("[Unfinished]"):
                                            paragraph.text = f"{i+1}. {paragraph.text.strip()}"
                                        else:
                                            paragraph.font.color.rgb = RGBColor(37, 133, 5)
                                            paragraph.text = f"{i+1}. {paragraph.text.strip()}"
                                    else:
                                        paragraph.text = f"{i+1}. {paragraph.text.strip()}"
                                else:
                                    paragraph.text = f"{i+1}. {paragraph.text.strip()}"
                            else:
                                paragraph.text = f"{i+1}. {paragraph.text.strip()}"
                    mid_sprint_key = str(app_name + "-midsprint")
                    if mid_sprint_key in new_planned_user_stories['iteration_data']:
                        text_frame = row.cells[2].text_frame
                        text_frame.clear()
                        text = rally_report_generator.process_tasks(new_planned_user_stories['iteration_data'][mid_sprint_key].split('\n'), team=mid_sprint_key)
                        text_frame.text = text if text else "None"
                        for i, paragraph in enumerate(text_frame.paragraphs):
                            paragraph.font.size = Pt(12)
                            if len(paragraph.text) > 4:
                                if paragraph.text.strip() not in rally_report_generator.get_stories_by_status(team=app_name, status="In-Progress"):
                                    if paragraph.text.strip() not in rally_report_generator.get_stories_by_status(team=app_name, status="Idea"):
                                        if paragraph.text.strip() not in rally_report_generator.get_stories_by_status(team=app_name, status="Defined"):
                                            if paragraph.text.strip().startswith("[Unfinished]"):
                                                paragraph.text = f"{i+1}. {paragraph.text.strip()}"
                                            else:
                                                paragraph.font.color.rgb = RGBColor(37, 133, 5)
                                                paragraph.text = f"{i+1}. {paragraph.text.strip()}"
                                        else:
                                            paragraph.text = f"{i+1}. {paragraph.text.strip()}"
                                    else:
                                        paragraph.text = f"{i+1}. {paragraph.text.strip()}"
                                else:
                                    paragraph.text = f"{i+1}. {paragraph.text.strip()}"
     
    elif len(table.rows[0].cells) == 2:
        if table.rows[0].cells[0].text.strip() == "Application" and table.rows[0].cells[1].text.strip().lower() == "planned user stories":
            for row in table.rows:
                app_name = row.cells[0].text.strip()
                if app_name in new_planned_user_stories['next_iteration_data']:
                    text_frame = row.cells[1].text_frame
                    text = new_planned_user_stories['next_iteration_data'][app_name]
                    text_frame.text = text if text else "None"
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
        elif table.rows[0].cells[0].text.strip() == "Application" and table.rows[0].cells[1].text.strip().lower() == "implemented user stories":
            for row in table.rows:
                app_name = row.cells[0].text.strip()
                if app_name in new_planned_user_stories['deployed_data']:
                    text_frame = row.cells[1].text_frame
                    titles = get_formatted_titles(new_planned_user_stories['deployed_data'], team=app_name)
                    text_frame.text = titles if titles else "None"
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
                else:
                    if row.cells[1].text.lower() != "implemented user stories":
                        text_frame = row.cells[1].text_frame
                        text_frame.text = "None"
                        text_frame.paragraphs[0].font.size = Pt(12)
        elif table.rows[0].cells[0].text.strip().lower() == "resource" and table.rows[0].cells[1].text.strip().lower() == "tasks":
            for row in table.rows:
                app_name = row.cells[0].text.strip()
                if any(key.startswith(app_name) for key in new_planned_user_stories['flex_resource_data']):
                    text_frame = row.cells[1].text_frame
                    text_frame.clear()
                    titles = new_planned_user_stories['flex_resource_data'][app_name]
                    text_frame.text = titles if titles else "None"
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
        elif table.rows[0].cells[0].text.strip() == "Completed Milestones" and table.rows[0].cells[1].text.strip() == "Incomplete Milestones":
            text_frame_active = table.rows[1].cells[0].text_frame
            text_frame_active.clear()
            
            milestone_counter = 1

            active_milestones_list = new_planned_user_stories['milestone_data'].get('Active', [])
            if active_milestones_list and rally_report_generator.milestone_update_counter < len(active_milestones_list):
                active_milestones = active_milestones_list[rally_report_generator.milestone_update_counter]
                for i in active_milestones:
                    milestone_title = f"{milestone_counter} {i.get('Milestone','')}"
                    p = text_frame_active.add_paragraph()
                    p.text = milestone_title
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(37, 133, 5)
                    for story in i.get('us', []):
                        story_titile = f"    • {story} "
                        p = text_frame_active.add_paragraph()
                        p.text = story_titile
                        p.font.size = Pt(13)
                        p.font.color.rgb = RGBColor(37, 133, 5)
                    p = text_frame_active.add_paragraph()
                    milestone_counter += 1

            text_frame_inactive = table.rows[1].cells[1].text_frame
            text_frame_inactive.clear()
            milestone_counter = 1

            inactive_milestones_list = new_planned_user_stories['milestone_data'].get('Inactive', [])
            if inactive_milestones_list and rally_report_generator.milestone_update_counter < len(inactive_milestones_list):
                inactive_milestones = inactive_milestones_list[rally_report_generator.milestone_update_counter]
                for i in inactive_milestones:
                    milestone_title = f"{milestone_counter} {i.get('Milestone','')}"
                    p = text_frame_inactive.add_paragraph()
                    p.text = milestone_title
                    p.font.size = Pt(14)
                    for story in i.get('us', []):
                        story_titile = f"    • {story} "
                        p = text_frame_inactive.add_paragraph()
                        p.text = story_titile
                        p.font.size = Pt(13)
                    p = text_frame_inactive.add_paragraph()
                    milestone_counter += 1
            rally_report_generator.milestone_update_counter += 1
    
    elif len(table.rows[0].cells) == 1:
        if table.rows[0].cells[0].text.strip() == "Milestone Report":
            text_frame = table.rows[1].cells[0].text_frame
            text_frame.clear()
            # Same Counter for Both Types.
            active_milestones_list = new_planned_user_stories['milestone_data'].get('Active', [])
            inactive_milestones_list = new_planned_user_stories['milestone_data'].get('Inactive', [])
            
            if active_milestones_list and rally_report_generator.milestone_update_counter < len(active_milestones_list):
                active_milestones = active_milestones_list[rally_report_generator.milestone_update_counter]
                for i in active_milestones:
                    milestone_title = f"{rally_report_generator.milestone_counter}. {i.get('Milestone','')}"
                    p = text_frame.add_paragraph()
                    p.text = milestone_title
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(37, 133, 5)
                    for story in i.get('us', []):
                        story_titile = f"    • {story} "
                        p = text_frame.add_paragraph()
                        p.text = story_titile
                        p.font.size = Pt(13)
                        p.font.color.rgb = RGBColor(37, 133, 5)
                    p = text_frame.add_paragraph()
                    rally_report_generator.milestone_counter += 1
                rally_report_generator.milestone_update_counter += 1

            elif inactive_milestones_list:
                inactive_milestones = inactive_milestones_list[rally_report_generator.milestone_update_counter - len(active_milestones_list)]
                for i in inactive_milestones:
                    milestone_title = f"{rally_report_generator.milestone_counter}. {i.get('Milestone','')}"
                    p = text_frame.add_paragraph()
                    p.text = milestone_title
                    p.font.size = Pt(14)
                    for story in i.get('us', []):
                        story_titile = f"    • {story} "
                        p = text_frame.add_paragraph()
                        p.text = story_titile
                        p.font.size = Pt(13)
                    p = text_frame.add_paragraph()
                    rally_report_generator.milestone_counter += 1
                rally_report_generator.milestone_update_counter += 1
            
            else:
                milestone_title = "None"
                p = text_frame.add_paragraph()
                p.text = milestone_title
                p.font.size = Pt(14)
                # Not sure
                rally_report_generator.milestone_update_counter += 1
                
            
def delete_all_except_first_row(table):    
    total_rows = len(table.rows)
    for i in range(total_rows - 1, 0, -1):
        table._tbl.remove(table._tbl.tr_lst[i])
   
def update_ppt_with_data(ppt_file: str, data: Dict, output_file: str, rally_report_generator: RallyReportGenerator, start_date, end_date: str = None, teams=[]) -> bool:
    try:
        prs = Presentation(ppt_file)
        if not end_date:
            slides=slides_to_remove(teams=teams)
            remove_slides(prs, slides)
            logging.info(f"Updating PPT {ppt_file} with fetched data.")
            print("test-1")

        # Find the index of milestone slide
        milestone_slides_count = len(data['milestone_data'].get('Active')) + len(data['milestone_data'].get('Inactive'))
        milestone_slide_index = -1
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    try:
                        first_row_cells = [cell.text.strip() for cell in shape.table.rows[0].cells]
                        if first_row_cells[0] == "Milestone Report":
                            milestone_slide_index = i
                            break  
                    except Exception:
                        continue  # Table might be empty or malformed
            if milestone_slide_index != -1:
                break  # Found it, break outer loop

        if milestone_slide_index != -1:
            original_slide = prs.slides[milestone_slide_index]
            for i in range(milestone_slides_count - 1):
                dest = prs.slides.add_slide(original_slide.slide_layout)
                for shp in original_slide.shapes:
                    el = shp.element
                    newel = copy.deepcopy(el)
                    dest.shapes._spTree.insert_element_before(newel, 'p:extLst')
                Slides = prs.slides._sldIdLst
                slides_to_move = Slides[len(Slides) - 1]
                Slides.insert(milestone_slide_index + i, slides_to_move)

        # dates in first slide 
        text_frame = prs.slides[0].shapes[2].text_frame
        text_frame.clear()
        iter_dates = rally_report_generator.fetch_iteration_dates(start_date=start_date)
        text_data = f"{iter_dates['iteration']} {iter_dates['Start']} to {iter_dates['end']}" if not end_date else f"{iter_dates['Start']} to {iter_dates['end']}"
        text_frame.text = text_data
        for paragraph in text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.size = Pt(16)  
        # Update teams in first slide 
        text_frame_1 = prs.slides[0].shapes[3].text_frame
        text_frame_2 = prs.slides[0].shapes[4].text_frame
        text_frame_1.clear()
        text_frame_2.clear()
    
        new_list = []
        for item in teams:
            if item == "ADB":
                new_list.extend(["EMIS", "CBIP", "D&B"])         
            elif item == "Data BAU":
                new_list.extend(["EDW", "Trade Credit", "KPI"])
            elif item == "SCUP NA Datamart - KPI":
                new_list.extend(["SCUP NA"])
            else:
                new_list.extend([item])

        new_list = sorted(new_list)

     # Determine which items go where
        items_for_frame1 = []
        items_for_frame2 = []      
        if len(new_list) > 5:
            items_for_frame1 = new_list[0:5]
            items_for_frame2 = new_list[5:]
        else:
            items_for_frame1 = new_list
    
        if len(text_frame_1.paragraphs) > 0 and text_frame_1.paragraphs[0].text == "":
            p = text_frame_1.paragraphs[0]._element
            p.getparent().remove(p)      
            for item in items_for_frame1:
                p = text_frame_1.add_paragraph()
                p.text = item
                p.level = 0  
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.size = Pt(16)    

        if items_for_frame2: 
            if len(text_frame_2.paragraphs) > 0 and text_frame_2.paragraphs[0].text == "":
                p = text_frame_2.paragraphs[0]._element
                p.getparent().remove(p)         
                for item in items_for_frame2:
                    p = text_frame_2.add_paragraph()
                    p.text = item
                    p.level = 0  # Apply first-level bullet (Removed paragraph.bullet = False)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.size = Pt(16)
        # updating all the tables
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    _update_table(shape.table, data, rally_report_generator)

        prs.save(output_file)
        logging.info(f"Updated PPT saved to {output_file}")
        return True
    except Exception as e:
        logging.error(f"Error updating PPT: {e}")
        return False
 

def slides_to_remove(teams):
    if teams == ["ADB", "RDM", "CDL", "CDH", "SCUP NA Datamart - KPI", "Data BAU"]:
        return []
    sot = {
        'ADB' : [1, 2, 8, 13],
        'RDM' : [3, 9, 14],
        'Data BAU' : [4, 10, 15],
        'SCUP NA Datamart - KPI' : [6, 12, 17]
    }

    lisT = []

    # For CDL and CDH 
    if (("CDL" not in teams) and ("CDH" not in teams)):
        lisT.extend([5, 11, 16])

    # For Everything else
    for team in sot.keys():
        if team not in teams:
            lisT.extend(sot[team])
    
    print(list(set(lisT)))
            
    return list(set(lisT))
    
def remove_slides(ppt, slides_indices):
    xml_slides = ppt.slides._sldIdLst
    slides = list(xml_slides)

    for index in sorted(slides_indices, reverse=True):
        xml_slides.remove(slides[index])

def validate_date_format(date_string: str) -> bool:
    try:
        datetime.strptime(date_string, '%Y-%m-%d')
        return True
    except ValueError:
        return False

def create_word_doc(milestone_details: Dict, output_file: str):
    if not milestone_details:
        return
    doc = Document()
    active_milestones_list = milestone_details['milestone_data'].get('Active', [])
    milestone_counter = 1
    if active_milestones_list :
        for i in active_milestones_list:
            milestone_title = f"{milestone_counter}. {i.get('Milestone','')}"
            milestone_counter += 1
            green_paragraph = doc.add_paragraph()
            green_run = green_paragraph.add_run(milestone_title)
            green_run.font.color.rgb = doc_color(37, 133, 5)

            for story in i.get('us', []):
                doc.add_paragraph(style='ListBullet',text= story)
                

    inactive_milestones_list = milestone_details['milestone_data'].get('Inactive', [])
    if inactive_milestones_list:
        for i in inactive_milestones_list:
            milestone_title = f"{milestone_counter}. {i.get('Milestone','')}"
            milestone_counter += 1
            blue_paragraph = doc.add_paragraph()
            blue_run = blue_paragraph.add_run(milestone_title)
            blue_run.font.color.rgb = doc_color(0,0,124)
    
            for story in i.get('us', []):
                doc.add_paragraph(style='ListBullet',text= story)

    doc.save(output_file)
    logging.info(f"Updated PPT saved to {output_file}")
    return True
    
def distribute_paragraphs_to_slides(paragraph_lengths, lines_per_slide=25):
    paragraph_indices = sorted(range(len(paragraph_lengths)), key=lambda i: paragraph_lengths[i], reverse=True)
    slides = []
    remaining_lines = []

    for index in paragraph_indices:
        length = paragraph_lengths[index]
        placed = False
        for i in range(len(slides)):
            if remaining_lines[i] >= length:
                slides[i].append(index)
                remaining_lines[i] -= length
                placed = True
                break
        if not placed:
            slides.append([index])
            remaining_lines.append(lines_per_slide - length)

    return slides

st.set_page_config(
    page_title="Rally Report Generator",
    page_icon="https://www.aigconnect.aig/Fallback/Assets/favicon.ico",
    layout="wide"
)
# Main UI
def main():
    st.title("Rally Report Generator")
    tracker = RallyReportGenerator()
    st.markdown("Select a template, date, and report type to generate and download your Rally report.")

    report_type_options = ["Sprint Report", "Milestone Report"]
    report_type = st.selectbox(
        label="Select Report Type",
        options=report_type_options,
        index=0,
        help="Choose the type of report"
    )
    with st.form(key="input_form"):
        selected_date_label = "Select Date" if report_type == 'Sprint Report' else "Select Start Date"
        selected_date = st.date_input(
            label=selected_date_label,
            value=datetime.today(),
            min_value=datetime(2020, 1, 1),
            max_value=datetime(2030, 12, 31),
            help="Choose any date that falls within the start and end dates of the identified iteration." if report_type == 'Sprint Report' else "Choose a date for the report"
        )
        if report_type == 'Milestone Report':
            selected_end_date = st.date_input(
                label="Select End Date",
                value=datetime.today(),
                min_value=datetime(2020, 1, 1),
                max_value=datetime(2030, 12, 31),
                help="Choose a date for the report"
            )
        else:
            selected_end_date = None

        # get projects based on there access?

        project_names = tracker.get_rally_projects()
        project_options = ["ALL", "ADB", "RDM", "CDL", "CDH", "SCUP NA Datamart - KPI", "Data BAU"]
        project_options_updated = []

        for project in project_options:
            if project in project_names:
                project_options_updated.append(project)
        project_options_updated.append("ALL")
        project_options_updated[:] = sorted(list(set(project_options_updated)))
        project_options_updated = ["SCUP NA" if (item == "SCUP NA Datamart - KPI")  else item for item in project_options_updated]



        template = st.multiselect(
            label="Select Teams",
            options=project_options_updated,
            default=project_options_updated[1],
            help="Choose one or more templates"
        )
        submit_button = st.form_submit_button(label="Generate Report")

    if submit_button:
        status_placeholder = st.empty()
        result_placeholder = st.empty()
        b64 = ''

        with st.spinner("Generating your Rally report... Please wait."):
            status_placeholder.info("Initializing report generator...")
            tracker = RallyReportGenerator()
            start_date = selected_date.strftime('%Y-%m-%d')
            
            
            status_placeholder.info("Fetching data from Rally...")

            if "ALL" in template:
                template = ["ADB", "RDM", "CDL", "CDH", "SCUP NA Datamart - KPI", "Data BAU"]
            
            template = ["Data BAU" if (item == "KPI" or item == "EDW" or item == "Trade Credit") else item for item in template]
            template = ["SCUP NA Datamart - KPI" if (item == "SCUP NA")  else item for item in template]

            # Report Type 
            if report_type == 'Sprint Report':
                end_date = None
                all_data = tracker.get_all_data(start_date=start_date, teams=template)
                print(all_data)
            else:
                end_date = selected_end_date.strftime('%Y-%m-%d')
                all_data = tracker.get_all_data(start_date=start_date, end_date=end_date, teams=template)
                

            if not all_data:
                status_placeholder.error("Data fetch failed. Cannot generate report.")
                return

            print(all_data)
            ppt_template = "template-sprintend.pptx" if report_type == 'Sprint Report' else "template-milestones.pptx"
            iteration = all_data['iter_dates'].get('iteration') + '_' if report_type == 'Sprint Report' else ''
            start = all_data['iter_dates'].get('Start')  if report_type == 'Sprint Report' else start_date
            end = all_data['iter_dates'].get('end') if report_type == 'Sprint Report' else end_date
            file = 'pptx' if report_type == 'Sprint Report' else 'docx'
            output_file = f"AIG_BAU_{report_type.replace(' ', '_')}_L3_{iteration}{start}_to_{end}.{file}"
            status_placeholder.info("Updating report with data...")

            success = update_ppt_with_data(
                ppt_file=ppt_template,
                data=all_data,
                teams=template,
                output_file=output_file,
                rally_report_generator=tracker,
                start_date=start_date,
                end_date = end_date) if not end_date else create_word_doc(milestone_details=all_data,output_file=output_file)

            if success:
                status_placeholder.success(f"Report generated successfully!")
                with open(output_file, "rb") as file:
                    pptx_data = file.read()
                b64 = BytesIO(pptx_data)
            else:
                status_placeholder.error("Error creating report. Please check the logs for details.")

            result_placeholder.download_button(
                    label="Download Report",
                    data=b64,
                    file_name=output_file,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation" if report_type == 'Sprint Report' else "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )
            os.remove(output_file)

if __name__ == "__main__":
    main()
