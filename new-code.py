from io import BytesIO
import streamlit as st
from datetime import datetime, timedelta
import logging
import os
import sys
import requests
from collections import defaultdict
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Pt
from typing import Dict, List, Optional, Tuple, Any
from dotenv import load_dotenv # Import dotenv

# Load environment variables from a .env file if it exists (for local development)
load_dotenv()

# Rely solely on environment variable for API key
API_KEY = os.environ.get("RALLY_API_KEY")
PPT_TEMPLATE_PATH = os.environ.get("PPT_TEMPLATE_PATH", "template-Mark.pptx") # Keep fallback for template path
LOG_LEVEL = os.environ.get("LOG_LEVEL", "INFO").upper()

# Define constants for frequently used values
TEAM_OWNERS_EMIS = ["lakshminarayana nainaru", "Govindarajan M", "Ravi Ranjan"]
STATUS_IN_PROGRESS = "In-Progress"
STATUS_IDEA = "Idea"
STATUS_DEFINED = "Defined"
STATUS_DEPLOYED = "Deployed"
COLOR_MIDSPRINT = '#fce205'
COLOR_MILESTONE_INACTIVE = "#21a2e0"

# Project IDs (replace with actual IDs if needed, or fetch dynamically)
PROJECT_IDS = {
    'RDM': 370746842872,
    'ADB': 81259836048,
    'CDH': 501753244176,
    'SCUP NA Datamart - KPI': 812858540333,
    'Data BAU': 343403251580
}

# Check if API key is set
if not API_KEY:
    logging.error("RALLY_API_KEY environment variable not set. Please set the API_KEY environment variable.")
    st.error("Rally API key not configured. Please set the RALLY_API_KEY environment variable.")
    sys.exit(1) # Exit if API key is not configured

# --- Logging ---
logging.basicConfig(level=LOG_LEVEL, format="%(asctime)s - %(levelname)s - %(message)s")

# --- RallyReportGenerator class and related functions ---
class RallyReportGenerator:
    """
    Fetches data from Rally and formats it for a report.

    Attributes:
        teams (Dict): Stores stories organized by team and status.
        status_index (Dict): Stores stories organized by status and team.
        headers (Dict): HTTP headers for Rally API requests.
    """
    def __init__(self, api_key: str = API_KEY):
        self.teams: Dict[str, Dict[str, Dict[str, Any]]] = {}
        # Note: self.status_index could potentially be derived from self.teams
        # if needed, but keeping both for now based on original code structure.
        self.status_index: Dict[str, Dict[str, Dict[str, Any]]] = {}
        self.headers: Dict[str, str] = {"zsessionid": api_key}
        logging.debug("RallyReportGenerator initialized.")

    def _make_request(self, url: str, params: Optional[Dict] = None) -> Optional[Dict]:
        """
        Makes an HTTP GET request to the Rally API.

        Args:
            url (str): The API endpoint URL.
            params (Optional[Dict]): Optional query parameters.

        Returns:
            Optional[Dict]: The JSON response data if successful, None otherwise.
        """
        try:
            response = requests.get(url, headers=self.headers, params=params)
            response.raise_for_status() # Raises HTTPError for bad responses (4xx or 5xx)
            return response.json()
        except requests.exceptions.HTTPError as e:
            logging.error(f"HTTP error: {e} for URL: {url} - Status Code: {e.response.status_code}")
            return None
        except requests.exceptions.RequestException as e:
            logging.error(f"Request error: {e} for URL: {url}")
            return None
        except Exception as e:
            logging.error(f"An unexpected error occurred: {e} for URL: {url}")
            return None

    def add_story(self, team_name: str, story_id: str, story_details: Dict) -> None:
        """
        Adds a story to the internal teams and status index dictionaries.

        Args:
            team_name (str): The name of the team the story belongs to.
            story_id (str): The unique ID of the story.
            story_details (Dict): A dictionary containing story details (must include 'status').
        """
        status = story_details.get('status')
        if not status:
            logging.warning(f"Story {story_id} for team {team_name} is missing status information. Skipping.")
            return

        self.teams.setdefault(team_name, {}).setdefault(status, {})[story_id] = story_details
        self.status_index.setdefault(status, {}).setdefault(team_name, {})[story_id] = story_details
        logging.debug(f"Story added: {team_name} - {story_id} - {status}")

    @st.cache_data(ttl=3600) # Cache for 1 hour
    def get_flex_resource_info(self, team: Optional[List[str]] = None, start_date: Optional[str] = None) -> Dict[str, str]:
        """
        Fetches flex resource information for a given team and date.

        Args:
            team (Optional[List[str]]): List of team member email addresses.
            start_date (Optional[str]): The start date for the iteration query (YYYY-MM-DD).

        Returns:
            Dict[str, str]: A dictionary mapping owner names to their tasks.
        """
        # Default team members if none provided
        team = team or ["megha.chakraborty@aig.com", "Aarthi.Panneerselvam@aig.com", "vukyam.srisravya@aig.com"]
        start_date = start_date or datetime.today().strftime('%Y-%m-%d')

        # Query for iteration based on start date
        url = f"https://rally1.rallydev.com/slm/webservice/v2.0/iteration?query=((StartDate <= \"{start_date}\") and (EndDate >= \"{start_date}\"))"
        flex_data: Dict[str, str] = defaultdict(str)
        counters: Dict[str, int] = defaultdict(int)

        response_data = self._make_request(url)

        if not response_data or not response_data.get('QueryResult', {}).get('Results'):
            logging.warning(f"No iteration found for flex resource info on {start_date}")
            return flex_data # Return empty dict if no iteration found

        # Assuming the first result is the relevant iteration
        iteration_ref = response_data['QueryResult']['Results'][0]['_ref']

        for owner in team:
            # Query for work products owned by the team member within the iteration
            query_params = {"query": f'(Owner = "{owner}")'}
            workproducts_url = f"{iteration_ref}/workproducts"
            workproducts = self._make_request(workproducts_url, params=query_params)

            if workproducts and workproducts.get('QueryResult', {}).get('Results'):
                for item in workproducts['QueryResult']['Results']:
                    name = item.get('Name')
                    # bau_team = item.get('c_DataBAUTeam', '') or item.get('Project', {}).get('_refObjectName') # Not used in flex_data
                    owner_name = item.get('Owner', {}).get('_refObjectName', 'Unknown Owner')
                    task_estimate = item.get('TaskEstimateTotal', 'N/A')
                    counters[owner_name] += 1
                    flex_data[owner_name] += f"{counters[owner_name]}. {name} ({task_estimate})\n"
            else:
                logging.debug(f"No work product found for owner {owner} in the iteration.")

        return flex_data

    @st.cache_data(ttl=3600) # Cache for 1 hour
    def fetch_iteration_dates(self, start_date: Optional[str] = None, next: bool = False) -> Any:
        """
        Fetches details for the current or next iteration based on a start date.

        Args:
            start_date (Optional[str]): A date within the target iteration (YYYY-MM-DD).
            next (bool): If True, fetches details for the iteration immediately following the found one.

        Returns:
            Any: A dictionary with iteration dates and name if successful,
                 or the start date of the next iteration if next=True, None otherwise.
        """
        start_date = start_date or datetime.today().strftime('%Y-%m-%d')

        # Find iteration based on start date
        url = f"https://rally1.rallydev.com/slm/webservice/v2.0/iteration?query=((StartDate <= \"{start_date}\") and (EndDate >= \"{start_date}\"))"
        response = self._make_request(url)

        if not response or not response.get('QueryResult', {}).get('Results'):
            logging.warning(f"No iteration found for date {start_date}")
            return None

        # Get details of the found iteration
        iteration_ref = response['QueryResult']['Results'][0]['_ref']
        response_data = self._make_request(iteration_ref)

        if not response_data or not response_data.get('Iteration'):
            logging.warning(f"Could not fetch details for iteration ref: {iteration_ref}")
            return None

        data = response_data['Iteration']

        if next:
            # Calculate the start date of the next iteration
            end_date_str = data.get('EndDate')
            if not end_date_str:
                logging.warning(f"Iteration {data.get('_refObjectName')} is missing EndDate.")
                return None
            end_date = datetime.strptime(end_date_str, "%Y-%m-%dT%H:%M:%S.%fZ")
            next_iteration_start_date = (end_date + timedelta(days=1)).strftime("%Y-%m-%d")
            return next_iteration_start_date
        else:
            # Return details of the current iteration
            start_date_str = data.get('StartDate')
            end_date_str = data.get('EndDate')
            if not start_date_str or not end_date_str:
                 logging.warning(f"Iteration {data.get('_refObjectName')} is missing StartDate or EndDate.")
                 return None

            start_date_formatted = datetime.strptime(start_date_str, "%Y-%m-%dT%H:%M:%S.%fZ").strftime('%Y-%m-%d')
            # Subtract one day from end date as per original logic
            end_date_formatted = (datetime.strptime(end_date_str, "%Y-%m-%dT%H:%M:%S.%fZ") - timedelta(days=1)).strftime('%Y-%m-%d')

            return {
                "Start": start_date_formatted,
                "end": end_date_formatted,
                "iteration": data.get('_refObjectName', 'Unknown Iteration'),
                "state": data.get('State', 'Unknown State')
            }

    def get_stories_by_status(self, status: Optional[str] = None, team: Optional[str] = None) -> Dict:
        """
        Retrieves stored stories filtered by status and/or team.

        Args:
            status (Optional[str]): The status to filter by.
            team (Optional[str]): The team to filter by.

        Returns:
            Dict: A dictionary containing the filtered stories.
        """
        if team and status:
            # Return list of titles for a specific team and status
            team_stories = self.teams.get(team, {})
            return [item['title'].strip() for item in team_stories.get(status, {}).values()]
        if team:
            # Return all stories for a specific team
            return self.teams.get(team, {})
        if status:
            # Return all stories for a specific status across all teams
            return self.status_index.get(status, {})
        # Return all stories across all teams and statuses
        return self.teams

    def _process_workproduct(self, project_name: str, workproduct: Dict, is_next_iter: Optional[bool] = False) -> Optional[Tuple[str, str]]:
        """
        Processes a single work product (User Story or Defect) and adds it to the report data.

        Args:
            project_name (str): The name of the Rally project the workproduct belongs to.
            workproduct (Dict): The dictionary containing workproduct details from Rally.
            is_next_iter (Optional[bool]): True if processing for the next iteration, False otherwise.

        Returns:
            Optional[Tuple[str, str]]: A tuple containing the data key and workproduct name if processed, None otherwise.
        """
        task_status = workproduct.get('ScheduleState', "")
        name = workproduct.get('Name', "")
        display_color = workproduct.get('DisplayColor', "")
        team = workproduct.get('c_DataBAUTeam', "") # Custom field for team
        owner = workproduct.get('Owner', {}).get('_refObjectName', "")

        imp_data_key = None
        story_key = f'{project_name}{len(self.teams.get(project_name, {}).get(task_status, {}))}' # Simple unique key

        # Determine the data key and add the story based on project, team, and color
        if project_name == "RDM":
            imp_data_key = 'RDM-midsprint' if display_color == COLOR_MIDSPRINT else 'RDM'
            if not is_next_iter:
                self.add_story('RDM', story_key, {'title': name, 'status': task_status})
        elif project_name == "Data BAU":
            if team == "KPI":
                imp_data_key = 'KPI-midsprint' if display_color == COLOR_MIDSPRINT else 'KPI'
                if not is_next_iter:
                    self.add_story('KPI', story_key, {'title': name, 'status': task_status})
            elif team == "EDW":
                imp_data_key = 'EDW-midsprint' if display_color == COLOR_MIDSPRINT else 'EDW'
                if not is_next_iter:
                    self.add_story('EDW', story_key, {'title': name, 'status': task_status})
            elif team == "Trade Credit":
                imp_data_key = 'Trade Credit-midsprint' if display_color == COLOR_MIDSPRINT else 'Trade Credit'
                if not is_next_iter:
                    self.add_story('Trade Credit', story_key, {'title': name, 'status': task_status})
        elif project_name == "CDL":
            imp_data_key = 'CDL-midsprint' if display_color == COLOR_MIDSPRINT else 'CDL'
            if not is_next_iter:
                self.add_story('CDL', story_key, {'title': name, 'status': task_status})
        elif project_name == "SCUP NA Datamart - KPI":
            imp_data_key = 'SCUP NA-midsprint' if display_color == COLOR_MIDSPRINT else 'SCUP NA'
            if not is_next_iter:
                self.add_story('SCUP NA', story_key, {'title': name, 'status': task_status})
        elif project_name == "CDH":
            imp_data_key = 'CDH-midsprint' if display_color == COLOR_MIDSPRINT else 'CDH'
            if not is_next_iter:
                self.add_story('CDH', story_key, {'title': name, 'status': task_status})
        elif project_name == "ADB":
            if display_color == COLOR_MIDSPRINT:
                if team == "EMIS":
                    if owner in TEAM_OWNERS_EMIS:
                        imp_data_key = 'EMIS Backend-midsprint'
                        if not is_next_iter:
                            self.add_story('EMIS Backend', f'EMIS Backend{len(self.teams.get("EMIS Backend", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                    else:
                        imp_data_key = 'EMIS UI-midsprint'
                        if not is_next_iter:
                            self.add_story('EMIS UI', f'EMIS UI{len(self.teams.get("EMIS UI", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                elif team == "CBIP":
                    imp_data_key = 'CBIP-midsprint'
                    if not is_next_iter:
                        self.add_story('CBIP', f'CBIP{len(self.teams.get("CBIP", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
            else: # Not mid-sprint color
                if team == "EMIS":
                    if owner in TEAM_OWNERS_EMIS:
                        imp_data_key = 'EMIS Backend'
                        if not is_next_iter:
                            self.add_story('EMIS Backend', f'EMIS Backend{len(self.teams.get("EMIS Backend", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                    else:
                        imp_data_key = 'EMIS UI'
                        if not is_next_iter:
                            self.add_story('EMIS UI', f'EMIS UI{len(self.teams.get("EMIS UI", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                elif team == "CBIP":
                    imp_data_key = 'CBIP'
                    if not is_next_iter:
                        self.add_story('CBIP', f'CBIP{len(self.teams.get("CBIP", {}).get(task_status, {}))}', {'title': name, 'status': task_status})

        # Return the determined key and name if a key was assigned
        return (imp_data_key, name) if imp_data_key else None

    @st.cache_data(ttl=3600) # Cache for 1 hour
    def fetch_iteration_data(self, start_date: Optional[str] = None, teams: List[str] = [], is_next_iter: Optional[bool] = False) -> Optional[Dict[str, str]]:
        """
        Fetches work product data for the iteration containing the start date, filtered by teams.

        Args:
            start_date (Optional[str]): A date within the target iteration (YYYY-MM-DD).
            teams (List[str]): A list of team names to filter by.
            is_next_iter (Optional[bool]): True if fetching for the next iteration, False otherwise.

        Returns:
            Optional[Dict[str, str]]: A dictionary with processed iteration data, None if fetch fails.
        """
        start_date = start_date or datetime.today().strftime('%Y-%m-%d')

        # Find iteration based on start date
        url = f"https://rally1.rallydev.com/slm/webservice/v2.0/iteration?query=((StartDate <= \"{start_date}\") and (EndDate >= \"{start_date}\"))"
        response_data = self._make_request(url)

        if not response_data or not response_data.get('QueryResult', {}).get('Results'):
            logging.warning(f"No iteration found for date {start_date} to fetch iteration data.")
            return None

        # Get details of the found iteration
        iter_details_ref = response_data['QueryResult']['Results'][0]['_ref']
        iter_details = self._make_request(iter_details_ref)

        if not iter_details or not iter_details.get('Iteration'):
            logging.warning(f"Could not fetch details for iteration ref: {iter_details_ref}")
            return None

        iter_data = iter_details['Iteration']
        project_name = iter_data.get('Project', {}).get('_refObjectName', 'Unknown Project')
        logging.debug(f"Processing iteration: {iter_data.get('Name', 'N/A')}, Project: {project_name}")

        # Initialize data dictionary with expected keys
        imp_data: Dict[str, str] = defaultdict(str, {
            'CBIP': '', 'EMIS Backend': '', 'EMIS UI': '',
            'CBIP-midsprint': '', 'EMIS Backend-midsprint': '',
            'EMIS UI-midsprint': '',
            'RDM': '', 'RDM-midsprint': '',
            'D&B': '', 'D&B-midsprint': '', # D&B might be under ADB project
            'EDW': '', 'EDW-midsprint': '', # EDW might be under Data BAU project
            'KPI': '', 'KPI-midsprint': '', # KPI might be under Data BAU or SCUP NA project
            'CDH': '', 'CDH-midsprint': '',
            'CDL': '', 'CDL-midsprint': '',
            'Trade Credit': '', 'Trade Credit-midsprint': '', # Trade Credit might be under Data BAU project
            'SCUP NA': '', 'SCUP NA-midsprint': '' # SCUP NA might be under SCUP NA Datamart - KPI project
        })
        counters: Dict[str, int] = defaultdict(int)

        # Check if the project of the iteration is in the selected teams
        if project_name in teams:
            workproducts_ref = iter_data.get('WorkProducts', {}).get('_ref', "")
            if workproducts_ref:
                workproducts_data = self._make_request(workproducts_ref)
                if workproducts_data and workproducts_data.get('QueryResult', {}).get('Results'):
                    workproducts = workproducts_data['QueryResult']['Results']
                    for workproduct in workproducts:
                        processed_result = self._process_workproduct(project_name, workproduct, is_next_iter)
                        if processed_result:
                            imp_data_key, name = processed_result
                            counters[imp_data_key] += 1
                            imp_data[imp_data_key] += f"{counters[imp_data_key]}. {name}\n"
                else:
                    logging.debug(f"No work products found for iteration {iter_data.get('Name', 'N/A')}.")
            else:
                 logging.debug(f"No work products reference found for iteration {iter_data.get('Name', 'N/A')}.")
        else:
            logging.debug(f"Iteration project '{project_name}' is not in the selected teams.")


        # The original code seems to fetch data for *all* iterations falling within the date range,
        # but only processes workproducts if the *iteration's project* is in the selected teams.
        # This might lead to missing data if workproducts for selected teams are in iterations
        # whose project is *not* in the selected teams.
        # A more robust approach might be to query workproducts directly filtered by project and iteration date range.
        # However, sticking closer to the original logic for now.

        # The original code iterates through response_data.get('QueryResult', {}).get('Results', [])
        # which implies processing multiple iterations if the date range overlaps.
        # Let's adjust to match that behavior more explicitly, although the initial query
        # `query=((StartDate <= \"{start_date}\") and (EndDate >= \"{start_date}\"))`
        # usually returns only one or zero iterations for a specific date.
        # If the intention is to get data for *all* iterations within a date range,
        # the loop structure is correct, but the filtering by `project_name in teams`
        # might need adjustment depending on the desired behavior.
        # For now, I'll keep the loop structure as in the original code, assuming it's intended
        # to handle cases where a date might span multiple iterations (less common).

        # Re-fetching response_data to iterate through results as in original code
        response_data_for_loop = self._make_request(url)
        if response_data_for_loop:
             for iteration in response_data_for_loop.get('QueryResult', {}).get('Results', []):
                iter_details = self._make_request(iteration.get('_ref', ""))
                if not iter_details:
                    continue

                iter_data = iter_details.get('Iteration', {})
                project_name = iter_data.get('Project', {}).get('_refObjectName', 'Unknown Project')

                # Only process iterations whose project is in the selected teams
                if project_name in teams:
                    workproducts_ref = iter_data.get('WorkProducts', {}).get('_ref', "")
                    if workproducts_ref:
                        workproducts_data = self._make_request(workproducts_ref)
                        if workproducts_data and workproducts_data.get('QueryResult', {}).get('Results'):
                            workproducts = workproducts_data['QueryResult']['Results']
                            for workproduct in workproducts:
                                # Pass the project_name to _process_workproduct
                                processed_result = self._process_workproduct(project_name, workproduct, is_next_iter)
                                if processed_result:
                                    imp_data_key, name = processed_result
                                    # Counters and imp_data are already being updated in _process_workproduct
                                    # Need to adjust _process_workproduct to return the data key and name
                                    # and update counters/imp_data here.
                                    # Let's refactor _process_workproduct to just return the key and name
                                    # and handle the adding/counting here.

                                    # Reverting _process_workproduct to return key and name, and handling add_story/counting here
                                    # This makes the logic clearer.
                                    # Let's redefine _process_workproduct to return (imp_data_key, name) or None
                                    # and move the add_story and counting logic here.

                                    # Re-implementing the logic from the original _process_workproduct here
                                    # based on the returned key and name.

                                    task_status = workproduct.get('ScheduleState', "")
                                    name = workproduct.get('Name', "")
                                    display_color = workproduct.get('DisplayColor', "")
                                    team_custom_field = workproduct.get('c_DataBAUTeam', "")
                                    owner = workproduct.get('Owner', {}).get('_refObjectName', "")

                                    current_imp_data_key = None
                                    story_key = f'{project_name}{len(self.teams.get(project_name, {}).get(task_status, {}))}' # Simple unique key

                                    if project_name == "RDM":
                                        current_imp_data_key = 'RDM-midsprint' if display_color == COLOR_MIDSPRINT else 'RDM'
                                        if not is_next_iter:
                                            self.add_story('RDM', story_key, {'title': name, 'status': task_status})
                                    elif project_name == "Data BAU":
                                        if team_custom_field == "KPI":
                                            current_imp_data_key = 'KPI-midsprint' if display_color == COLOR_MIDSPRINT else 'KPI'
                                            if not is_next_iter:
                                                self.add_story('KPI', story_key, {'title': name, 'status': task_status})
                                        elif team_custom_field == "EDW":
                                            current_imp_data_key = 'EDW-midsprint' if display_color == COLOR_MIDSPRINT else 'EDW'
                                            if not is_next_iter:
                                                self.add_story('EDW', story_key, {'title': name, 'status': task_status})
                                        elif team_custom_field == "Trade Credit":
                                            current_imp_data_key = 'Trade Credit-midsprint' if display_color == COLOR_MIDSPRINT else 'Trade Credit'
                                            if not is_next_iter:
                                                self.add_story('Trade Credit', story_key, {'title': name, 'status': task_status})
                                    elif project_name == "CDL":
                                        current_imp_data_key = 'CDL-midsprint' if display_color == COLOR_MIDSPRINT else 'CDL'
                                        if not is_next_iter:
                                            self.add_story('CDL', story_key, {'title': name, 'status': task_status})
                                    elif project_name == "SCUP NA Datamart - KPI":
                                        current_imp_data_key = 'SCUP NA-midsprint' if display_color == COLOR_MIDSPRINT else 'SCUP NA'
                                        if not is_next_iter:
                                            self.add_story('SCUP NA', story_key, {'title': name, 'status': task_status})
                                    elif project_name == "CDH":
                                        current_imp_data_key = 'CDH-midsprint' if display_color == COLOR_MIDSPRINT else 'CDH'
                                        if not is_next_iter:
                                            self.add_story('CDH', story_key, {'title': name, 'status': task_status})
                                    elif project_name == "ADB":
                                        if display_color == COLOR_MIDSPRINT:
                                            if team_custom_field == "EMIS":
                                                if owner in TEAM_OWNERS_EMIS:
                                                    current_imp_data_key = 'EMIS Backend-midsprint'
                                                    if not is_next_iter:
                                                        self.add_story('EMIS Backend', f'EMIS Backend{len(self.teams.get("EMIS Backend", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                                                else:
                                                    current_imp_data_key = 'EMIS UI-midsprint'
                                                    if not is_next_iter:
                                                        self.add_story('EMIS UI', f'EMIS UI{len(self.teams.get("EMIS UI", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                                            elif team_custom_field == "CBIP":
                                                current_imp_data_key = 'CBIP-midsprint'
                                                if not is_next_iter:
                                                    self.add_story('CBIP', f'CBIP{len(self.teams.get("CBIP", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                                        else: # Not mid-sprint color
                                            if team_custom_field == "EMIS":
                                                if owner in TEAM_OWNERS_EMIS:
                                                    current_imp_data_key = 'EMIS Backend'
                                                    if not is_next_iter:
                                                        self.add_story('EMIS Backend', f'EMIS Backend{len(self.teams.get("EMIS Backend", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                                                else:
                                                    current_imp_data_key = 'EMIS UI'
                                                    if not is_next_iter:
                                                        self.add_story('EMIS UI', f'EMIS UI{len(self.teams.get("EMIS UI", {}).get(task_status, {}))}', {'title': name, 'status': task_status})
                                            elif team_custom_field == "CBIP":
                                                current_imp_data_key = 'CBIP'
                                                if not is_next_iter:
                                                    self.add_story('CBIP', f'CBIP{len(self.teams.get("CBIP", {}).get(task_status, {}))}', {'title': name, 'status': task_status})

                                    if current_imp_data_key:
                                        counters[current_imp_data_key] += 1
                                        imp_data[current_imp_data_key] += f"{counters[current_imp_data_key]}. {name}\n"

                        else:
                            logging.debug(f"No work products found for iteration {iter_data.get('Name', 'N/A')}.")
                    else:
                         logging.debug(f"No work products reference found for iteration {iter_data.get('Name', 'N/A')}.")
                else:
                    logging.debug(f"Iteration project '{project_name}' is not in the selected teams.")


        return imp_data


    @st.cache_data(ttl=3600) # Cache for 1 hour
    def get_all_data(self, start_date: Optional[str] = None, teams: List[str] = [], end_date: Optional[str] = None) -> Dict[str, Any]:
        """
        Fetches all necessary data for generating a report.

        Args:
            start_date (Optional[str]): The start date for the report (YYYY-MM-DD).
            teams (List[str]): A list of team names to filter data by.
            end_date (Optional[str]): The end date for the report (YYYY-MM-DD).

        Returns:
            Dict[str, Any]: A dictionary containing all fetched and processed data.
        """
        start_date = start_date or datetime.today().strftime('%Y-%m-%d')

        # Fetch iteration dates first as they are needed for other calls
        iter_dates = self.fetch_iteration_dates(start_date=start_date)
        if not iter_dates:
            logging.error(f"Could not fetch iteration dates for {start_date}.")
            # Return partial data or None, depending on how critical iter_dates are
            # Returning partial data might allow generating a milestone report
            # even if sprint data fails.
            iter_dates = {} # Provide an empty dict to avoid errors later

        data: Dict[str, Any] = {
            "iter_dates": iter_dates,
            "milestone_data" : None # Initialize to None
        }

        if end_date:
            # Milestone report only
            data["milestone_data"] = self.milestonedetails(start_date=start_date, teams=teams, end_date=end_date)
        else:
            # Sprint report
            # Fetch data for the current iteration
            data["iteration_data"] = self.fetch_iteration_data(start_date=start_date, teams=teams, is_next_iter=False)
            # Fetch deployed stories (status="Deployed")
            data["deployed_data"] = self.get_stories_by_status(status=STATUS_DEPLOYED)
            # Fetch milestone data for the current iteration date range
            data["milestone_data"] = self.milestonedetails(start_date=start_date, teams=teams)
            # Fetch data for the next iteration
            next_iter_start_date = self.fetch_iteration_dates(start_date=start_date, next=True)
            if next_iter_start_date:
                 data["next_iteration_data"] = self.fetch_iteration_data(start_date=next_iter_start_date, teams=teams, is_next_iter=True)
            else:
                 logging.warning("Could not fetch start date for the next iteration.")
                 data["next_iteration_data"] = None # Ensure key exists even if data is missing

            # Flex resource data - keeping this commented out as in the original code
            # data["flex_resource_data"] = self.get_flex_resource_info(start_date=start_date)


        return data

    @st.cache_data(ttl=3600) # Cache for 1 hour
    def milestonedetails(self, start_date: Optional[str] = None, end_date: Optional[str] = None, teams: List[str] = [] ) -> Optional[Dict[str, List[Dict[str, Any]]]]:
        """
        Fetches milestone details within a specified date range for selected projects.

        Args:
            start_date (Optional[str]): The start date for the milestone query (YYYY-MM-DD).
            end_date (Optional[str]): The end date for the milestone query (YYYY-MM-DD).
            teams (List[str]): A list of team names to filter projects by.

        Returns:
            Optional[Dict[str, List[Dict[str, Any]]]]: A dictionary containing active and inactive milestones, None if fetch fails.
        """
        query_start_date = start_date
        query_end_date = end_date

        if not query_start_date or not query_end_date:
            # If dates are not provided, use current iteration dates
            iter_dates = self.fetch_iteration_dates(start_date=start_date)
            if not iter_dates:
                logging.error("Could not fetch iteration dates to determine milestone date range.")
                return None
            query_start_date = iter_dates.get('Start')
            query_end_date = iter_dates.get('end')

        if not query_start_date or not query_end_date:
             logging.error("Milestone start or end date is missing.")
             return None

        # Validate date formats
        if not validate_date_format(query_start_date):
            logging.error(f"Invalid start date format for milestone: {query_start_date}. Please use YYYY-MM-DD format.")
            # In a real app, you might raise an exception or return an error indicator
            return None

        if not validate_date_format(query_end_date):
            logging.error(f"Invalid end date format for milestone: {query_end_date}. Please use YYYY-MM-DD format.")
             # In a real app, you might raise an exception or return an error indicator
            return None

        # Filter project IDs based on selected teams
        filtered_project_ids = [
            project_id for team_name, project_id in PROJECT_IDS.items()
            if team_name in teams
        ]

        if not filtered_project_ids:
            logging.warning("No valid project IDs found for the selected teams to fetch milestones.")
            return defaultdict(list, {'Active': [], 'Inactive': []}) # Return empty structure

        milestone_data: Dict[str, List[Dict[str, Any]]] = defaultdict(list, {
            'Active': [],
            'Inactive': []
        })

        # Fetch milestones for each filtered project ID
        for project_id in filtered_project_ids:
            # Rally query to find milestones within the date range for the project
            url = (
                f'https://rally1.rallydev.com/slm/webservice/v2.0/milestone?'
                f'fetch=Artifacts,Name,DisplayColor,FormattedID,TargetDate&'
                f'query=(((Projects contains "/project/{project_id}") AND (TargetDate >= "{query_start_date}")) AND (TargetDate <= "{query_end_date}"))&'
                f'start=1&pagesize=25&order=TargetDate DESC'
            )
            response_data = self._make_request(url)

            if not response_data or not response_data.get('QueryResult', {}).get('Results'):
                logging.debug(f"No milestones found for project ID {project_id} within the date range.")
                continue # Continue to the next project

            for milestone in response_data['QueryResult']['Results']:
                milestone_name = milestone.get('_refObjectName','')
                milestone_id = milestone.get('FormattedID','')
                milestone_color = milestone.get('DisplayColor','')
                artifact_ref = milestone.get('Artifacts', {}).get('_ref','')

                artifact_data: List[str] = []
                if artifact_ref:
                    artifact_details = self._make_request(artifact_ref)
                    if artifact_details and artifact_details.get('QueryResult', {}).get('Results'):
                         artifact_data = [result.get('_refObjectName', 'Unknown Artifact') for result in artifact_details['QueryResult']['Results']]
                    else:
                        logging.debug(f"No artifacts found for milestone {milestone_id}.")


                milestone_entry = {"Milestone": f"{milestone_name} {milestone_id}", "color": milestone_color, "us": artifact_data}

                if milestone_color == COLOR_MILESTONE_INACTIVE:
                    milestone_data['Inactive'].append(milestone_entry)
                else:
                    milestone_data['Active'].append(milestone_entry)

        return milestone_data

    def process_tasks(self, tasks_list: List[str], team: str) -> str:
        """
        Formats and sorts a list of task strings for display.

        Args:
            tasks_list (List[str]): A list of task strings (e.g., ["1. Task Name"]).
            team (str): The name of the team the tasks belong to.

        Returns:
            str: A single string with processed and sorted task names.
        """
        def task_order_key(task_str: str) -> bool:
            """Helper function to determine sort order based on task status."""
            # Extract task name after the number and dot
            if len(task_str.split(". ", 1)) < 2:
                 return False # Cannot parse, treat as not in specified statuses
            task_name = task_str.split(". ", 1)[1].strip()

            # Return False if the task is in In-Progress, Idea, or Defined status
            # This will make these tasks appear later in the sorted list (False < True)
            return not (task_name in self.get_stories_by_status(team=team, status=STATUS_IN_PROGRESS) or
                        task_name in self.get_stories_by_status(team=team, status=STATUS_IDEA) or
                        task_name in self.get_stories_by_status(team=team, status=STATUS_DEFINED))

        # Filter out empty strings and sort based on the helper key
        sorted_tasks = sorted([task for task in tasks_list if task], key=task_order_key)

        # Extract just the task name (remove the leading number and dot)
        tasks_string = "\n".join(task.split(". ", 1)[1] for task in sorted_tasks if len(task.split(". ", 1)) > 1)

        return tasks_string

# --- Helper Functions ---

def get_formatted_titles(data: Dict, team: str) -> str:
    """
    Formats a dictionary of stories for a team into a numbered list string.

    Args:
        data (Dict): The dictionary containing story data (e.g., from get_stories_by_status).
        team (str): The name of the team.

    Returns:
        str: A string with formatted story titles.
    """
    # Assuming data structure is {team_name: {status: {story_id: {title: "...", ...}}}}
    # Need to get all titles for the team, regardless of status, if data is from get_stories_by_status(status="Deployed")
    # If data is from get_stories_by_status(status="Deployed"), it's already filtered by status.
    # The original code calls get_stories_by_status(status="Deployed") and passes the whole result here.
    # Let's assume data is the result of get_stories_by_status(status="Deployed"), which is {team_name: {story_id: {title: "...", ...}}}
    team_stories = data.get(team, {})
    titles = [value.get("title", "Unknown Title").strip() for value in team_stories.values()]
    return "\n".join(f"{i+1}. {title}" for i, title in enumerate(titles))


def _update_table(table: Any, data: Dict, rally_report_generator: RallyReportGenerator) -> None:
    """
    Updates a table shape in a PowerPoint slide with data.

    Args:
        table (Any): The python-pptx table shape object.
        data (Dict): The data dictionary containing report information.
        rally_report_generator (RallyReportGenerator): The report generator instance.
    """
    # Check the number of columns to determine table type
    num_cols = len(table.rows[0].cells)
    header_row_text = [cell.text.strip().lower() for cell in table.rows[0].cells]

    if num_cols > 2 and header_row_text[0] == "application" and header_row_text[1] == "planned user stories" and header_row_text[2] == "mid sprint user stories":
        # Update Planned and Mid Sprint User Stories table
        iteration_data = data.get('iteration_data', {})
        if not iteration_data:
            logging.warning("No iteration data available to update planned/mid sprint table.")
            return

        for row in table.rows[1:]: # Skip header row
            app_name = row.cells[0].text.strip()
            if app_name in iteration_data:
                # Update Planned User Stories column
                planned_text_frame = row.cells[1].text_frame
                planned_text_frame.clear()
                # Split the raw string data into a list of tasks
                planned_tasks_raw = iteration_data.get(app_name, "").split('\n')
                # Process and format the tasks
                planned_tasks_formatted = rally_report_generator.process_tasks(planned_tasks_raw, team=app_name)
                planned_text_frame.text = planned_tasks_formatted if planned_tasks_formatted else "None"
                # Apply formatting (size and color for completed tasks)
                for i, paragraph in enumerate(planned_text_frame.paragraphs):
                    paragraph.font.size = Pt(12)
                    # Check if the task name (after number and dot) is NOT in In-Progress, Idea, or Defined status
                    task_name = paragraph.text.split(". ", 1)[1].strip() if len(paragraph.text.split(". ", 1)) > 1 else paragraph.text.strip()
                    if task_name and task_name not in rally_report_generator.get_stories_by_status(team=app_name, status=STATUS_IN_PROGRESS) and \
                       task_name not in rally_report_generator.get_stories_by_status(team=app_name, status=STATUS_IDEA) and \
                       task_name not in rally_report_generator.get_stories_by_status(team=app_name, status=STATUS_DEFINED):
                        paragraph.font.color.rgb = RGBColor(37, 133, 5) # Green color for completed/accepted

                # Update Mid Sprint User Stories column
                mid_sprint_key = f"{app_name}-midsprint"
                if mid_sprint_key in iteration_data:
                    mid_sprint_text_frame = row.cells[2].text_frame
                    mid_sprint_text_frame.clear()
                    # Split the raw string data into a list of tasks
                    mid_sprint_tasks_raw = iteration_data.get(mid_sprint_key, "").split('\n')
                     # Process and format the tasks
                    mid_sprint_tasks_formatted = rally_report_generator.process_tasks(mid_sprint_tasks_raw, team=mid_sprint_key) # Use mid_sprint_key for team? Or app_name? Original uses mid_sprint_key. Let's stick to original.
                    mid_sprint_text_frame.text = mid_sprint_tasks_formatted if mid_sprint_tasks_formatted else "None"
                    # Apply formatting (size and color for completed tasks)
                    for i, paragraph in enumerate(mid_sprint_text_frame.paragraphs):
                         paragraph.font.size = Pt(12)
                         # Check if the task name (after number and dot) is NOT in In-Progress, Idea, or Defined status
                         task_name = paragraph.text.split(". ", 1)[1].strip() if len(paragraph.text.split(". ", 1)) > 1 else paragraph.text.strip()
                         # The original code checks against app_name statuses here, not mid_sprint_key.
                         # Let's stick to the original logic.
                         if task_name and task_name not in rally_report_generator.get_stories_by_status(team=app_name, status=STATUS_IN_PROGRESS) and \
                            task_name not in rally_report_generator.get_stories_by_status(team=app_name, status=STATUS_IDEA) and \
                            task_name not in rally_report_generator.get_stories_by_status(team=app_name, status=STATUS_DEFINED):
                             paragraph.font.color.rgb = RGBColor(37, 133, 5) # Green color for completed/accepted


    elif num_cols == 2:
        if header_row_text[0] == "application" and header_row_text[1] == "planned user stories":
            # Update Next Iteration Planned User Stories table
            next_iteration_data = data.get('next_iteration_data', {})
            if not next_iteration_data:
                 logging.warning("No next iteration data available to update planned stories table.")
                 return

            for row in table.rows[1:]: # Skip header row
                app_name = row.cells[0].text.strip()
                if app_name in next_iteration_data:
                    text_frame = row.cells[1].text_frame
                    text = next_iteration_data.get(app_name, "")
                    text_frame.text = text if text else "None"
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(12)

        elif header_row_text[0] == "application" and header_row_text[1] == "implemented user stories":
            # Update Implemented User Stories table
            deployed_data = data.get('deployed_data', {})
            if not deployed_data:
                 logging.warning("No deployed data available to update implemented stories table.")
                 # Clear existing content if no data
                 for row in table.rows[1:]:
                     text_frame = row.cells[1].text_frame
                     text_frame.clear()
                     text_frame.text = "None"
                     text_frame.paragraphs[0].font.size = Pt(12)
                 return

            for row in table.rows[1:]: # Skip header row
                app_name = row.cells[0].text.strip()
                # get_formatted_titles expects the structure {team_name: {story_id: {title: "...", ...}}}
                # deployed_data from get_stories_by_status(status="Deployed") is {team_name: {story_id: {title: "...", ...}}}
                titles = get_formatted_titles(deployed_data, team=app_name)
                text_frame = row.cells[1].text_frame
                text_frame.clear() # Clear existing content before adding new
                text_frame.text = titles if titles else "None"
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(12)

        elif header_row_text[0] == "resource" and header_row_text[1] == "tasks":
            # Update Flex Resource Tasks table
            flex_resource_data = data.get('flex_resource_data', {})
            if not flex_resource_data:
                 logging.warning("No flex resource data available to update tasks table.")
                 # Clear existing content if no data
                 for row in table.rows[1:]:
                     text_frame = row.cells[1].text_frame
                     text_frame.clear()
                     text_frame.text = "None"
                     text_frame.paragraphs[0].font.size = Pt(12)
                 return

            for row in table.rows[1:]: # Skip header row
                resource_name = row.cells[0].text.strip()
                # Check if the resource name (or a key starting with it) exists in flex_resource_data
                # The original code checks if any key starts with app_name (which is resource_name here)
                # Let's assume resource_name directly maps to a key in flex_resource_data
                if resource_name in flex_resource_data:
                     text_frame = row.cells[1].text_frame
                     text_frame.clear()
                     titles = flex_resource_data.get(resource_name, "")
                     text_frame.text = titles if titles else "None"
                     for paragraph in text_frame.paragraphs:
                         paragraph.font.size = Pt(12)
                else:
                    logging.debug(f"No flex resource data found for resource: {resource_name}")
                    # Optionally clear the cell if no data is found for the resource
                    text_frame = row.cells[1].text_frame
                    text_frame.clear()
                    text_frame.text = "None"
                    text_frame.paragraphs[0].font.size = Pt(12)


        elif header_row_text[0] == "completed milestones" and header_row_text[1] == "incomplete milestones":
            # Update Milestones table
            milestone_data = data.get('milestone_data', {})
            if not milestone_data:
                 logging.warning("No milestone data available to update milestones table.")
                 # Clear existing content if no data
                 for row in table.rows[1:]:
                     for cell in row.cells:
                         text_frame = cell.text_frame
                         text_frame.clear()
                         text_frame.text = "None"
                         text_frame.paragraphs[0].font.size = Pt(12)
                 return

            # Assuming the milestone table has only one data row below the header
            if len(table.rows) > 1:
                completed_cell = table.rows[1].cells[0]
                incomplete_cell = table.rows[1].cells[1]

                # Update Completed Milestones cell
                completed_text_frame = completed_cell.text_frame
                completed_text_frame.clear()
                milestone_counter = 1
                for milestone_entry in milestone_data.get('Active', []):
                    milestone_title = f"{milestone_counter}. {milestone_entry.get('Milestone','')}"
                    p = completed_text_frame.add_paragraph()
                    p.text = milestone_title
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(37, 133, 5) # Green color
                    for story in milestone_entry.get('us',[]):
                        story_title = f"     • {story} "
                        p = completed_text_frame.add_paragraph()
                        p.text = story_title
                        p.font.size = Pt(13)
                        p.font.color.rgb = RGBColor(37, 133, 5) # Green color
                    # Add a blank paragraph for spacing between milestones
                    if milestone_entry != milestone_data['Active'][-1]: # Don't add after the last one
                        completed_text_frame.add_paragraph()
                    milestone_counter += 1
                if not milestone_data.get('Active'):
                    completed_text_frame.text = "None"
                    completed_text_frame.paragraphs[0].font.size = Pt(12)


                # Update Incomplete Milestones cell
                incomplete_text_frame = incomplete_cell.text_frame
                incomplete_text_frame.clear()
                milestone_counter = 1
                for milestone_entry in milestone_data.get('Inactive', []):
                    milestone_title = f"{milestone_counter}. {milestone_entry.get('Milestone','')}"
                    p = incomplete_text_frame.add_paragraph()
                    p.text = milestone_title
                    p.font.size = Pt(14)
                    # No color change for incomplete milestones as per original code
                    for story in milestone_entry.get('us',[]):
                        story_title = f"     • {story} "
                        p = incomplete_text_frame.add_paragraph()
                        p.text = story_title
                        p.font.size = Pt(13)
                    # Add a blank paragraph for spacing between milestones
                    if milestone_entry != milestone_data['Inactive'][-1]: # Don't add after the last one
                         incomplete_text_frame.add_paragraph()
                    milestone_counter += 1
                if not milestone_data.get('Inactive'):
                    incomplete_text_frame.text = "None"
                    incomplete_text_frame.paragraphs[0].font.size = Pt(12)


def delete_all_except_first_row(table: Any) -> None:
    """
    Deletes all rows from a table except the header row.

    Args:
        table (Any): The python-pptx table shape object.
    """
    # Ensure table has more than one row before attempting deletion
    if len(table.rows) > 1:
        total_rows = len(table.rows)
        # Iterate backwards to avoid index issues
        for i in range(total_rows - 1, 0, -1):
            # Access the underlying XML element to remove the row
            table._tbl.remove(table._tbl.tr_lst[i])
    else:
        logging.debug("Table has only one row, no rows to delete.")


def update_ppt_with_data(ppt_file: str, data: Dict, output_file: str, rally_report_generator: RallyReportGenerator, start_date: str, end_date: Optional[str] = None, teams: List[str] = []) -> bool:
    """
    Updates a PowerPoint presentation with fetched Rally data.

    Args:
        ppt_file (str): The path to the input PowerPoint template file.
        data (Dict): The dictionary containing the data to populate the PPT.
        output_file (str): The path to save the updated PowerPoint file.
        rally_report_generator (RallyReportGenerator): The report generator instance.
        start_date (str): The start date used for fetching data (YYYY-MM-DD).
        end_date (Optional[str]): The end date used for fetching data (YYYY-MM-DD), for milestone reports.
        teams (List[str]): The list of teams selected for the report.

    Returns:
        bool: True if the update was successful, False otherwise.
    """
    try:
        prs = Presentation(ppt_file)

        if not end_date: # Logic specific to Sprint Report
            logging.info(f"Updating Sprint Report PPT {ppt_file} with fetched data.")

            # Determine which slides to remove based on selected teams
            slides_indices_to_remove = slides_to_remove(teams=teams)
            remove_slides(prs, slides_indices_to_remove)

            # Update dates in the first slide (assuming slide 0 is the title slide)
            # Find the shapes containing the date and team information.
            # This part is fragile as it relies on the specific structure of the template.
            # A more robust approach would be to use placeholders or specific text markers.
            date_shape = None
            team_shape_1 = None
            team_shape_2 = None

            # Iterate through shapes to find potential date and team shapes
            for shape in prs.slides[0].shapes:
                 # Look for shapes that might contain the date text
                 if shape.has_text_frame:
                     text = shape.text_frame.text.strip()
                     # Look for patterns that might indicate the date shape
                     if "iteration" in text.lower() or (" to " in text and any(char.isdigit() for char in text)):
                         date_shape = shape
                     # Look for shapes that might contain team names (assuming bullet points or similar)
                     # This is a heuristic and might need adjustment based on the template
                     if len(shape.text_frame.paragraphs) > 1 and any(p.text.strip() for p in shape.text_frame.paragraphs):
                         if team_shape_1 is None:
                             team_shape_1 = shape
                         elif team_shape_2 is None:
                             team_shape_2 = shape

            if date_shape:
                text_frame = date_shape.text_frame
                text_frame.clear()
                iter_dates = data.get('iter_dates', {}) # Use data dict
                if iter_dates:
                    text_frame.text = f"{iter_dates.get('iteration', 'N/A')} {iter_dates.get('Start', 'N/A')} to {iter_dates.get('end', 'N/A')}"
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255) # White color
                        paragraph.font.size = Pt(16)
                else:
                    text_frame.text = "Iteration Dates Not Available"
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(16)

            else:
                logging.warning("Could not find the date shape on the first slide.")


            # Update teams in the first slide
            # This logic assumes two text boxes for teams.
            # Need to map selected teams to the names displayed on the slide (e.g., "ADB" maps to "EMIS", "CBIP", "D&B")
            displayed_teams = []
            for item in teams:
                if item == "ADB":
                    displayed_teams.extend(["EMIS", "CBIP", "D&B"])
                elif item == "Data BAU":
                    displayed_teams.extend(["EDW", "Trade Credit", "KPI"])
                elif item == "SCUP NA Datamart - KPI":
                    displayed_teams.extend(["SCUP NA"])
                else:
                    # Assume other selected teams are displayed as is
                    displayed_teams.extend([item])

            displayed_teams = sorted(list(set(displayed_teams))) # Remove duplicates and sort

            if team_shape_1:
                text_frame_1 = team_shape_1.text_frame
                text_frame_1.clear()
                if len(displayed_teams) > 5:
                    text_frame_1.text  = "\n".join(displayed_teams[0:5])
                    for paragraph in text_frame_1.paragraphs:
                        paragraph.level = 0 # Ensure no bullet points if not intended
                        paragraph.font.color.rgb = RGBColor(255, 255, 255) # White color
                        paragraph.font.size = Pt(16)
                else:
                    text_frame_1.text =  "\n".join(displayed_teams)
                    for paragraph in text_frame_1.paragraphs:
                        paragraph.level = 0
                        paragraph.font.color.rgb = RGBColor(255, 255, 255) # White color
                        paragraph.font.size = Pt(16)

            else:
                 logging.warning("Could not find the first team shape on the first slide.")


            if team_shape_2:
                text_frame_2 = team_shape_2.text_frame
                text_frame_2.clear()
                if len(displayed_teams) > 5:
                    text_frame_2.text = "\n".join(displayed_teams[5::])
                    for paragraph in text_frame_2.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255) # White color
                        paragraph.font.size = Pt(16)
                        paragraph.bullet = False # Ensure no bullet points
                # If less than or equal to 5 teams, the second shape might be empty or used for something else.
                # Clear it if it was intended for teams and there are no teams for it.
                elif len(displayed_teams) <= 5 and team_shape_2.text_frame.text.strip():
                     text_frame_2.clear() # Clear if it had previous content
                     text_frame_2.text = "" # Set text to empty string

            else:
                 logging.warning("Could not find the second team shape on the first slide.")


            # updating all the tables in Sprint Report slides
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_table:
                        _update_table(shape.table, data, rally_report_generator)

        else: # Logic specific to Milestone Report
            logging.info(f"Updating Milestone Report PPT {ppt_file} with fetched data.")

            # For Milestone Report, the original code only updates the milestone table.
            # Find the slide containing the milestone table.
            milestone_slide = None
            for slide in prs.slides:
                 for shape in slide.shapes:
                     if shape.has_table:
                         # Check if the table headers match the milestone table headers
                         if len(shape.table.rows) > 0 and len(shape.table.rows[0].cells) == 2:
                             header_text = [cell.text.strip().lower() for cell in shape.table.rows[0].cells]
                             if header_text[0] == "completed milestones" and header_text[1] == "incomplete milestones":
                                 milestone_slide = slide
                                 break # Found the milestone slide
                 if milestone_slide:
                     break # Exit outer loop once slide is found

            if milestone_slide:
                for shape in milestone_slide.shapes:
                    if shape.has_table:
                        # Assuming there's only one table on the milestone slide with the correct headers
                        if len(shape.table.rows) > 0 and len(shape.table.rows[0].cells) == 2:
                             header_text = [cell.text.strip().lower() for cell in shape.table.rows[0].cells]
                             if header_text[0] == "completed milestones" and header_text[1] == "incomplete milestones":
                                _update_table(shape.table, data, rally_report_generator)
                                break # Updated the milestone table
            else:
                logging.warning("Could not find the milestone table in the presentation.")


        prs.save(output_file)
        logging.info(f"Updated PPT saved to {output_file}")
        return True
    except FileNotFoundError:
        logging.error(f"PowerPoint template file not found at {ppt_file}")
        return False
    except Exception as e:
        logging.error(f"Error updating PPT: {e}", exc_info=True) # Log traceback
        return False

def slides_to_remove(teams: List[str]) -> List[int]:
    """
    Determines which slides to remove from the Sprint Report template based on selected teams.

    Args:
        teams (List[str]): The list of selected team names.

    Returns:
        List[int]: A list of slide indices to remove.
    """
    # Mapping of team names to slide indices to remove if the team is *not* selected
    # This mapping is based on the structure of the original template.
    # Indices are 0-based.
    slide_mapping: Dict[str, List[int]] = {
        'ADB' : [1, 2, 8], # Example: If ADB is NOT selected, remove slides 1, 2, and 8
        'RDM' : [3, 9],
        'Data BAU' : [4, 10, 14],
        # CDL and CDH seem to share slides 5 and 11 based on the original logic
        # 'CDL' : [5, 11],
        # 'CDH' : [5, 11],
        'SCUP NA Datamart - KPI' : [6, 12]
    }

    indices_to_remove = []

    # Handle CDL and CDH shared slides
    if "CDL" not in teams and "CDH" not in teams:
         indices_to_remove.extend([5, 11])

    # Handle other teams based on the mapping
    # Note: The input teams list might contain "KPI", "EDW", "Trade Credit", "SCUP NA"
    # which map to "Data BAU" and "SCUP NA Datamart - KPI" projects in Rally.
    # The logic here should use the project names that correspond to the slides.
    # Let's assume the keys in slide_mapping are the project names relevant to the slides.
    # We need to check if any of the selected teams map to these project names.

    # Create a set of relevant project names from the selected teams
    relevant_projects = set()
    for team in teams:
        if team in ["ADB", "RDM", "CDL", "CDH"]:
            relevant_projects.add(team)
        elif team in ["KPI", "EDW", "Trade Credit"]:
            relevant_projects.add("Data BAU")
        elif team == "SCUP NA":
            relevant_projects.add("SCUP NA Datamart - KPI")
        # Add other direct team-to-project mappings if necessary

    # Iterate through the slide mapping keys (project names)
    for project_name, indices in slide_mapping.items():
        # If the project name is NOT in the set of relevant projects, add its indices to the removal list
        if project_name not in relevant_projects:
            indices_to_remove.extend(indices)

    # Handle additional removal logic from the original code
    # This logic seems dependent on combinations of other slides being removed.
    # It might need careful review based on the actual template structure.
    # Assuming 13 and 15 are also team-specific summary slides.
    if (1 in indices_to_remove and 3 in indices_to_remove): # If ADB and RDM slides are removed
        indices_to_remove.append(13)
    if (5 in indices_to_remove and 6 in indices_to_remove): # If CDL/CDH and SCUP NA slides are removed
        indices_to_remove.append(15)


    # Return unique sorted indices
    return sorted(list(set(indices_to_remove)))

def remove_slides(ppt: Presentation, slides_indices: List[int]) -> None:
    """
    Removes slides from a PowerPoint presentation by index.

    Args:
        ppt (Presentation): The python-pptx Presentation object.
        slides_indices (List[int]): A list of 0-based indices of slides to remove.
    """
    # Ensure indices are valid and within the range of slides
    valid_indices = [i for i in slides_indices if 0 <= i < len(ppt.slides)]
    if len(valid_indices) != len(slides_indices):
        logging.warning(f"Some requested slide indices were out of bounds. Valid indices: {valid_indices}")

    # Remove slides in reverse order to avoid index shifts
    xml_slides = ppt.slides._sldIdLst
    slides = list(xml_slides) # Create a list from the slide IDs

    for index in sorted(valid_indices, reverse=True):
        logging.debug(f"Removing slide at index: {index}")
        try:
            xml_slides.remove(slides[index])
        except IndexError:
            logging.error(f"Failed to remove slide at index {index}. Index might be invalid after previous removals.")


def validate_date_format(date_string: str) -> bool:
    """
    Validates if a string is in YYYY-MM-DD format.

    Args:
        date_string (str): The string to validate.

    Returns:
        bool: True if the format is valid, False otherwise.
    """
    try:
        datetime.strptime(date_string, '%Y-%m-%d')
        return True
    except ValueError:
        return False

# --- Streamlit App ---

st.set_page_config(
    page_title="Rally Report Generator",
    page_icon="https://www.aigconnect.aig/Fallback/Assets/favicon.ico",
    layout="wide"
)

def main():
    """Main function to run the Streamlit application."""
    st.title("Rally Report Generator")
    st.markdown("Select a template, date, and report type to generate and download your Rally report.")

    report_type_options = ["Sprint Report", "Milestone Report"]
    report_type = st.selectbox(
        label="Select Report Type",
        options=report_type_options,
        index=0,
        help="Choose the type of report"
    )

    with st.form(key="input_form"):
        selected_date_label = "Select Date (within iteration)" if report_type == 'Sprint Report' else "Select Start Date"
        selected_date = st.date_input(
            label=selected_date_label,
            value=datetime.today(),
            min_value=datetime(2020, 1, 1),
            max_value=datetime(2030, 12, 31),
            help="Choose any date that falls within the start and end dates of the identified iteration." if report_type == 'Sprint Report' else "Choose a start date for the report"
        )

        selected_end_date = None
        if report_type == 'Milestone Report':
            selected_end_date = st.date_input(
                label="Select End Date",
                value=datetime.today(),
                min_value=datetime(2020, 1, 1),
                max_value=datetime(2030, 12, 31),
                help="Choose an end date for the report"
            )

        # Team selection options - map user-friendly names to Rally project names where necessary
        teams_options_display = ["ALL", "ADB", "RDM", "CDL", "CDH", "SCUP NA", "KPI", "EDW", "Trade Credit"]
        selected_teams_display = st.multiselect(
            label="Select Teams",
            options=teams_options_display,
            default=teams_options_display[0],
            help="Choose one or more teams to include in the report."
        )

        submit_button = st.form_submit_button(label="Generate Report")

    if submit_button:
        # Map display team names to internal project names used in the script
        selected_teams_internal = []
        if "ALL" in selected_teams_display:
             # Use the project names that correspond to the slides/data fetching logic
             selected_teams_internal = ["ADB", "RDM", "CDL", "CDH", "SCUP NA Datamart - KPI", "Data BAU"]
        else:
            for team_display in selected_teams_display:
                if team_display == "SCUP NA":
                    selected_teams_internal.append("SCUP NA Datamart - KPI")
                elif team_display in ["KPI", "EDW", "Trade Credit"]:
                     selected_teams_internal.append("Data BAU")
                else:
                    # Assume other display names match internal project names
                    selected_teams_internal.append(team_display)

        # Remove duplicates and ensure selected_teams_internal only contains valid project names
        valid_project_names = list(PROJECT_IDS.keys()) + ["Data BAU", "SCUP NA Datamart - KPI"] # Add Data BAU and SCUP NA Datamart - KPI which are used in logic but not in PROJECT_IDS directly
        selected_teams_internal = list(set([team for team in selected_teams_internal if team in valid_project_names]))

        if not selected_teams_internal:
             st.warning("Please select at least one valid team.")
             return


        status_placeholder = st.empty()
        result_placeholder = st.empty()
        b64 = None # Initialize b64 to None

        with st.spinner("Generating your Rally report... Please wait."):
            status_placeholder.info("Initializing report generator...")
            tracker = RallyReportGenerator()
            start_date_str = selected_date.strftime('%Y-%m-%d')

            status_placeholder.info("Fetching data from Rally...")

            # Fetch data based on report type
            all_data = tracker.get_all_data(
                start_date=start_date_str,
                teams=selected_teams_internal,
                end_date=selected_end_date.strftime('%Y-%m-%d') if selected_end_date else None
            )

            if not all_data or (report_type == 'Sprint Report' and not all_data.get('iteration_data')):
                status_placeholder.error("Data fetch failed or returned no relevant data. Cannot generate report. Check logs for details.")
                return

            # Determine template and output file name
            ppt_template = "template.pptx" if report_type == 'Sprint Report' else "Mtemplate.pptx"

            # Construct output file name based on report type and dates
            output_filename_parts = ["AIG_BAU", report_type.replace(' ', '_'), "L3"]
            iter_info = all_data.get('iter_dates', {})
            if report_type == 'Sprint Report' and iter_info:
                 output_filename_parts.append(f"{iter_info.get('iteration', 'UnknownIteration')}_")
                 output_filename_parts.append(f"{iter_info.get('Start', start_date_str)}_to_{iter_info.get('end', 'UnknownEnd')}")
            elif report_type == 'Milestone Report':
                 output_filename_parts.append(f"{start_date_str}_to_{selected_end_date.strftime('%Y-%m-%d') if selected_end_date else 'UnknownEnd'}")
            else:
                 output_filename_parts.append(f"{start_date_str}_to_UnknownEnd") # Fallback filename


            output_file = "_".join(output_filename_parts) + ".pptx"


            status_placeholder.info("Updating PowerPoint with data...")

            success = update_ppt_with_data(
                ppt_file=ppt_template,
                data=all_data,
                teams=selected_teams_internal,
                output_file=output_file,
                rally_report_generator=tracker,
                start_date=start_date_str,
                end_date=selected_end_date.strftime('%Y-%m-%d') if selected_end_date else None
            )

            if success:
                status_placeholder.success(f"Report generated successfully!")
                try:
                    with open(output_file, "rb") as file:
                        pptx_data = file.read()
                    b64 = BytesIO(pptx_data)
                except FileNotFoundError:
                    status_placeholder.error(f"Generated file not found: {output_file}")
                    logging.error(f"Generated file not found: {output_file}")
                    b64 = None # Ensure b64 is None if file not found
            else:
                status_placeholder.error("Failed to update PowerPoint. Please check the logs for details.")

            # Provide download button only if b64 data is available
            if b64:
                result_placeholder.download_button(
                        label="Download Report",
                        data=b64,
                        file_name=output_file,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
            else:
                 result_placeholder.warning("Report file could not be generated or found for download.")

            # Clean up the generated file
            if os.path.exists(output_file):
                try:
                    os.remove(output_file)
                    logging.debug(f"Removed temporary file: {output_file}")
                except OSError as e:
                    logging.warning(f"Could not remove temporary file {output_file}: {e}")


if __name__ == "__main__":
    main()
