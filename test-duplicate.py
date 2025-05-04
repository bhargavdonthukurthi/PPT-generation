from pptx import Presentation
import copy

def duplicate_slide(prs, index):
    """Duplicate the slide with the given index in prs.

    Args:
        prs (pptx.Presentation): The presentation object.
        index (int): The zero-based index of the slide to duplicate.

    Returns:
        pptx.slide.Slide: The newly created duplicated slide.
    """
    source = prs.slides[index]
    try:
        blank_slide_layout = prs.slide_layouts[7]  # Assuming a blank layout exists
    except IndexError:
        blank_slide_layout = prs.slide_layouts[-1] # Or take the last one if no blank

    dest = prs.slides.add_slide(blank_slide_layout)

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    # for key, value in source.rels.items():
    #     if not "notesSlide" in value.reltype:
    #         dest.rels.add_relationship(value.reltype, value._target, value.rId)

    return dest

if __name__ == '__main__':
    prs = Presentation("template.pptx") # Load the presentation with 2 slides

    slide_to_duplicate_index = 7  # Index of the second slide (0-based)
    num_duplicates = 1  # Specify the number of times you want to duplicate

    for _ in range(num_duplicates):
        duplicate_slide(prs, slide_to_duplicate_index)

    prs.save("your_presentation_with_multiple_duplicates.pptx")
    print(f"Slide at index {slide_to_duplicate_index} duplicated {num_duplicates} times!")